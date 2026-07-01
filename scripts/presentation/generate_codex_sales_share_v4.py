from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from generate_codex_sales_share_v2 import (
    COLORS,
    SLIDE_H,
    SLIDE_W,
    background,
    browser_mock,
    card,
    footer,
    header,
    line,
    rect,
    set_shape_text,
    small_card,
    step_circle,
    text_box,
)


OUT_DIR = Path(r"D:\ad-ops-workbench\outputs\codex_sales_share_v4")
PPTX_PATH = OUT_DIR / "sales_codex_quick_start_v4_for_colleagues.pptx"
QUICK_CARD = OUT_DIR / "sales_codex_one_page_quick_start.md"


def page_title(slide, page, title, subtitle=None):
    header(slide, page, title, subtitle, None)


def prompt_box(slide, x, y, w, h, title, prompt):
    box = rect(slide, x, y, w, h, fill="111827", line=None, radius=True)
    set_shape_text(
        box,
        title,
        prompt,
        title_size=15,
        body_size=13.2,
        title_color="FFFFFF",
        body_color="E5E7EB",
        margin=0.22,
    )
    return box


def cover(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    rect(slide, 0.78, 0.76, 11.78, 5.66, fill="111827", line=None, radius=True)
    text_box(slide, "销售版 Codex", 1.35, 1.28, 10.5, 0.55, size=34, color="FFFFFF", bold=True, align="center")
    text_box(slide, "快速上手卡", 1.35, 1.98, 10.5, 0.5, size=28, color="93C5FD", bold=True, align="center")
    text_box(slide, "适合刚安装 Codex 的销售同事：先装技能包，再打开页面，复制一句话开始只读分析。", 2.0, 2.85, 9.35, 0.42, size=16, color="E5E7EB", align="center")
    labels = [("1 装技能包", COLORS["blue"]), ("2 打开页面", COLORS["orange"]), ("3 复制一句话", COLORS["green"])]
    for i, (txt, color) in enumerate(labels):
        box = rect(slide, 2.05 + i * 3.2, 4.05, 2.35, 0.86, fill="FFFFFF", line=None, radius=True)
        set_shape_text(box, txt, None, title_size=15.5, title_color=color, title_bold=True, align="center", valign="middle", margin=0.04)
    text_box(slide, "先只读，不直接修改任何内容。", 1.35, 5.55, 10.5, 0.25, size=13, color="CBD5E1", align="center")
    footer(slide, 1)


def three_steps(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    page_title(slide, 2, "第一次使用，只做这三步", "不要先研究复杂功能。先让 Codex 在一个你熟悉的页面上跑一遍只读分析。")
    steps = [
        ("1", "确认已安装技能包", "技能包文件夹：amazon-sales-starter"),
        ("2", "打开一个业务页面", "广告、库存、Amazon 前台、竞品搜索结果都可以"),
        ("3", "复制第一句测试", "只读帮我看当前页面，不要直接修改"),
    ]
    colors = [COLORS["blue"], COLORS["orange"], COLORS["green"]]
    for i, (num, title, body) in enumerate(steps):
        x = 1.0 + i * 4.05
        box = rect(slide, x, 2.0, 3.35, 3.45, fill="FFFFFF", line="D7DEE9", radius=True)
        step_circle(slide, x + 0.3, 2.38, num, colors[i])
        text_box(slide, title, x + 0.3, 3.15, 2.5, 0.3, size=17, color=COLORS["ink"], bold=True)
        text_box(slide, body, x + 0.3, 3.78, 2.55, 0.6, size=13.5, color=COLORS["muted"])
        if i < 2:
            line(slide, x + 3.52, 3.68, x + 3.84, 3.68, color="CBD5E1", width=1.5)
    rect(slide, 2.0, 6.0, 9.35, 0.56, fill="111827", line=None, radius=True)
    text_box(slide, "成功标准：Codex 能说清楚它看到了什么、哪些判断有证据、哪些地方还不确定。", 2.28, 6.16, 8.8, 0.16, size=12.5, color="FFFFFF", bold=True, align="center")


def why_skill(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    page_title(slide, 3, "为什么要先装销售技能包？", "刚下载的 Codex 不知道你的销售职责，也不知道 SKU、ACOS、滞销这些业务词。")
    card(slide, 0.95, 2.0, 5.25, 3.65, "不装技能包", "你要反复解释：\n什么是 SKU / ASIN\n广告字段怎么看\n哪些动作不能乱改\n销售判断先看什么", accent=COLORS["red"])
    card(slide, 7.08, 2.0, 5.25, 3.65, "装好技能包", "它先知道：\n销售常见职责\n常用术语\n页面读取顺序\n只读和动作边界", accent=COLORS["green"])
    text_box(slide, "不是让你带新人，而是让一个读过销售手册的助手先跑第一遍。", 1.8, 6.13, 9.8, 0.28, size=14.5, color=COLORS["ink"], bold=True, align="center")


def install_skill(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    page_title(slide, 4, "技能包放哪里？", "安装公司版 Codex 后，把销售技能包文件夹放进本机 skills 目录。")
    left = rect(slide, 0.95, 1.85, 5.45, 3.8, fill="FFFFFF", line="D7DEE9", radius=True)
    set_shape_text(left, "技能包文件夹", "amazon-sales-starter\n\n里面应该能看到：\nSKILL.md\nreferences 文件夹\nagents 文件夹", title_size=18, body_size=15, title_color=COLORS["blue"], body_color="334155", margin=0.28)
    right = rect(slide, 7.0, 1.85, 5.45, 3.8, fill="111827", line=None, radius=True)
    set_shape_text(right, "常见 Windows 路径", "C:\\Users\\<用户名>\\.codex\\skills\\\n\n放好后：\n...\\.codex\\skills\\amazon-sales-starter\\SKILL.md", title_size=18, body_size=14.5, title_color="FFFFFF", body_color="E5E7EB", margin=0.28)
    text_box(slide, "如果找不到目录，直接找助理远程协助，不要自己硬研究。", 2.0, 6.1, 9.35, 0.28, size=13.5, color=COLORS["orange"], bold=True, align="center")


def first_prompt(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    page_title(slide, 5, "第一句测试：直接复制这一段", "打开任意一个你平时看的业务页面，然后把下面这段发给 Codex。")
    prompt_box(
        slide,
        1.35,
        1.95,
        10.65,
        2.6,
        "复制给 Codex",
        "只读帮我看当前页面。\n按 Amazon 销售入门技能包，先识别页面类型、可见字段、初步判断、风险和下一步建议。\n不要直接修改任何内容。",
    )
    rows = [
        ("它应该先说", "这是什么页面、看到了哪些对象和字段"),
        ("它不应该做", "不应该直接改广告、改价格、改 listing"),
        ("你要检查", "结论后面有没有证据，是否有不确定项"),
    ]
    for i, (title, body) in enumerate(rows):
        small_card(slide, 1.35 + i * 3.55, 5.05, 3.0, 0.86, title, body, [COLORS["blue_soft"], COLORS["red_soft"], COLORS["green_soft"]][i], [COLORS["blue"], COLORS["red"], COLORS["green"]][i])


def expected_output(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    page_title(slide, 6, "它第一次输出，应该长这样", "不要只看它写得像不像人话，要看它有没有证据、有无风险、有无下一步。")
    sections = [
        ("看到的信息", "页面类型、SKU/ASIN/搜索词、可见字段"),
        ("初步判断", "当前问题更像曝光、点击、转化还是效率"),
        ("支持证据", "每个判断后面对应哪个字段或页面信息"),
        ("风险/不确定", "哪些字段看不懂、哪些结论样本不足"),
        ("下一步建议", "先补证据，还是生成待确认动作清单"),
    ]
    for i, (title, body) in enumerate(sections):
        x = 0.95 + (i % 2) * 6.05
        y = 1.82 + (i // 2) * 1.25
        if i == 4:
            x = 0.95
            w = 11.28
        else:
            w = 5.25
        small_card(slide, x, y, w, 0.82, title, body, [COLORS["blue_soft"], COLORS["teal_soft"], COLORS["orange_soft"], COLORS["red_soft"], COLORS["green_soft"]][i], [COLORS["blue"], COLORS["teal"], COLORS["orange"], COLORS["red"], COLORS["green"]][i])
    rect(slide, 2.15, 6.0, 9.05, 0.55, fill="111827", line=None, radius=True)
    text_box(slide, "如果它只给一句“建议优化”，就让它补：证据在哪里？反证是什么？", 2.45, 6.15, 8.45, 0.16, size=12.5, color="FFFFFF", bold=True, align="center")


def current_page(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    page_title(slide, 7, "最推荐从“当前页面只读”开始", "不需要导表，不需要截图十几张。先让它看你已经打开的页面。")
    browser_mock(slide, 0.95, 1.82, 5.6, 3.75, "已打开 Chrome 页面", ["广告页面", "库存/产品页面", "Amazon 前台页面", "竞品搜索结果"], active_idx=1)
    right = rect(slide, 7.05, 1.82, 5.2, 3.75, fill="FFFFFF", line="D7DEE9", radius=True)
    set_shape_text(right, "你要做", "打开页面\n确认自己能看到数据\n复制只读提示词\n检查它输出的证据", title_size=18, body_size=16, title_color=COLORS["green"], body_color="334155", margin=0.3)
    text_box(slide, "第一阶段只追求：它帮你少看一遍、少整理一遍。", 2.0, 6.12, 9.35, 0.28, size=14, color=COLORS["ink"], bold=True, align="center")


def competitor_card(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    page_title(slide, 8, "用法 1：竞品市场快速整理", "适合：你要看一个搜索词、一组竞品、一个新品机会。")
    prompt_box(slide, 0.95, 1.85, 5.35, 2.1, "复制提示词", "帮我看这个搜索词下的 10 个竞品，整理价格、评论、主图和卖点差异。最后告诉我我们能不能打、还缺什么证据。")
    right = rect(slide, 7.0, 1.85, 5.35, 2.1, fill="FFFFFF", line="D7DEE9", radius=True)
    set_shape_text(right, "你会拿到", "价格带\n评论门槛\n主图打法\n卖点结构\n差距和机会点", title_size=16, body_size=13.8, title_color=COLORS["blue"], body_color="334155", margin=0.22)
    checks = [
        ("检查 1", "是否只看了单个链接？如果是，让它扩成 10 个竞品。"),
        ("检查 2", "是否有“我们能不能打”的结论？"),
        ("检查 3", "是否标出还缺什么证据？"),
    ]
    for i, (t, b) in enumerate(checks):
        small_card(slide, 0.95 + i * 4.05, 4.8, 3.35, 0.9, t, b, COLORS["blue_soft"], COLORS["blue"])


def conversion_card(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    page_title(slide, 9, "用法 2：点击后不转化，先找证据", "适合：有点击、有花费，但订单少，想判断是不是 listing 或竞争力问题。")
    prompt_box(slide, 0.95, 1.85, 5.35, 2.05, "复制提示词", "我怀疑这个 SKU 不是没流量，而是点击后不转化。请帮我找支持证据、反证、还缺什么证据，以及下一步最小检查动作。")
    items = [
        ("支持证据", "点击、CVR、价格、评价、竞品对照"),
        ("反证", "样本太少、流量词不准、时间窗不对"),
        ("下一步", "先补页面/竞品/广告哪类证据"),
    ]
    for i, (t, b) in enumerate(items):
        small_card(slide, 7.0, 1.85 + i * 1.0, 5.35, 0.78, t, b, [COLORS["green_soft"], COLORS["red_soft"], COLORS["orange_soft"]][i], [COLORS["green"], COLORS["red"], COLORS["orange"]][i])
    rect(slide, 2.0, 5.85, 9.35, 0.55, fill="111827", line=None, radius=True)
    text_box(slide, "关键：不是让它替你拍脑袋，而是让它先把证据找齐。", 2.3, 6.0, 8.75, 0.16, size=12.5, color="FFFFFF", bold=True, align="center")


def ad_card(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    page_title(slide, 10, "用法 3：广告页面只读归类", "适合：你已经打开广告页面，想先分清问题属于哪一类。")
    browser_mock(slide, 0.95, 1.82, 5.4, 3.65, "广告页面字段", ["曝光 / 点击", "花费 / CPC", "订单 / 销售", "CTR / CVR / ACOS"], active_idx=3)
    prompt_box(slide, 7.0, 1.82, 5.35, 1.78, "复制提示词", "只读看这个广告页面，先判断问题更像曝光不足、点击不足、点击后不转化，还是效率问题。每个判断后面写证据。")
    right = rect(slide, 7.0, 4.0, 5.35, 1.45, fill="FFFFFF", line="D7DEE9", radius=True)
    set_shape_text(right, "不要一上来让它调广告", "先让它归类问题，再生成待确认动作清单。", title_size=15, body_size=12.5, title_color=COLORS["red"], body_color="334155", margin=0.2)


def action_card(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    page_title(slide, 11, "用法 4：动作清单，不是直接执行", "适合：你想让 Codex 帮你整理可操作建议，但还没准备让它改系统。")
    prompt_box(slide, 0.95, 1.85, 5.35, 2.05, "复制提示词", "按规则生成待确认动作清单，不要直接执行。清单里写清楚：对象、建议动作、依据、风险、执行后怎么回读验证。")
    table = rect(slide, 7.0, 1.85, 5.35, 3.55, fill="FFFFFF", line="D7DEE9", radius=True)
    set_shape_text(table, "清单必须包含", "对象\n建议动作\n依据\n风险\n回读验证方式", title_size=17, body_size=15, title_color=COLORS["blue"], body_color="334155", margin=0.28)
    rect(slide, 2.0, 5.92, 9.35, 0.55, fill="111827", line=None, radius=True)
    text_box(slide, "看到清单后，人确认，再谈执行。", 2.3, 6.07, 8.75, 0.16, size=12.8, color="FFFFFF", bold=True, align="center")


def glossary(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    page_title(slide, 12, "这些基础词，技能包已经先教给它", "所以你不用每次从 SKU、ASIN、ACOS 开始解释。")
    terms = [
        ("SKU / ASIN", "产品身份"),
        ("CTR / CVR", "点击和转化"),
        ("ACOS / CPC", "广告效率"),
        ("FBA / 库存天数", "库存状态"),
        ("滞销", "库存周转风险"),
        ("空运 / 海运利润", "利润口径"),
    ]
    fills = [COLORS["blue_soft"], COLORS["teal_soft"], COLORS["orange_soft"], COLORS["green_soft"], COLORS["red_soft"], COLORS["purple_soft"]]
    accents = [COLORS["blue"], COLORS["teal"], COLORS["orange"], COLORS["green"], COLORS["red"], COLORS["purple"]]
    for i, (t, b) in enumerate(terms):
        x = 0.95 + (i % 3) * 4.05
        y = 1.9 + (i // 3) * 1.45
        small_card(slide, x, y, 3.35, 0.98, t, b, fills[i], accents[i])
    text_box(slide, "公司内部更细的规则，它不知道就应该标“待确认”，不要装懂。", 2.05, 5.95, 9.25, 0.28, size=13.5, color=COLORS["ink"], bold=True, align="center")


def repair_prompts(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    page_title(slide, 13, "如果它答得不好，直接这样追问", "不要重新来一遍。用追问把它拉回证据和边界。")
    prompts = [
        ("太空泛", "把每个结论后面的证据字段列出来。"),
        ("乱下结论", "把支持证据和反证分开写。"),
        ("想直接改", "先生成待确认动作清单，不要执行。"),
        ("看不懂字段", "把你不理解的业务词列出来，我逐个确认。"),
        ("没有下一步", "给我一个最小检查动作，不要给大而全建议。"),
        ("缺少边界", "重写一次，明确哪些结论只基于当前页面。"),
    ]
    for i, (t, b) in enumerate(prompts):
        x = 0.88 + (i % 2) * 6.1
        y = 1.82 + (i // 2) * 1.22
        small_card(slide, x, y, 5.3, 0.88, t, b, COLORS["orange_soft"], COLORS["orange"])


def boundaries(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    page_title(slide, 14, "四条边界，先记住", "能不能放心用，关键看边界是否清楚。")
    items = [
        ("只读优先", "新手先不让它修改系统"),
        ("不越权", "只看你有权限或你提供的内容"),
        ("不迷信", "结论要能回到证据核对"),
        ("不装懂", "内部规则不确定就标待确认"),
    ]
    for i, (t, b) in enumerate(items):
        x = 0.95 + (i % 2) * 6.08
        y = 2.0 + (i // 2) * 1.55
        small_card(slide, x, y, 5.35, 1.05, t, b, [COLORS["blue_soft"], COLORS["red_soft"], COLORS["green_soft"], COLORS["purple_soft"]][i], [COLORS["blue"], COLORS["red"], COLORS["green"], COLORS["purple"]][i])
    rect(slide, 2.2, 5.75, 8.95, 0.62, fill="111827", line=None, radius=True)
    text_box(slide, "让它先跑第一遍，不是让它替你担责任。", 2.5, 5.93, 8.35, 0.16, size=13, color="FFFFFF", bold=True, align="center")


def one_page_summary(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    page_title(slide, 15, "一页速查：从这里开始", "把这页当成你的第一天使用卡。")
    left = rect(slide, 0.95, 1.8, 5.35, 4.25, fill="FFFFFF", line="D7DEE9", radius=True)
    set_shape_text(left, "使用顺序", "1. 确认 amazon-sales-starter 已安装\n2. 打开一个业务页面\n3. 复制只读提示词\n4. 看它是否列出证据和风险\n5. 需要动作时，先要待确认清单", title_size=18, body_size=14.5, title_color=COLORS["blue"], body_color="334155", margin=0.28)
    right = rect(slide, 7.0, 1.8, 5.35, 4.25, fill="111827", line=None, radius=True)
    set_shape_text(right, "第一句", "只读帮我看当前页面。\n按 Amazon 销售入门技能包，先识别页面类型、可见字段、初步判断、风险和下一步建议。\n不要直接修改任何内容。", title_size=18, body_size=14, title_color="FFFFFF", body_color="E5E7EB", margin=0.28)
    footer(slide, 15)


SLIDES = [
    cover,
    three_steps,
    why_skill,
    install_skill,
    first_prompt,
    expected_output,
    current_page,
    competitor_card,
    conversion_card,
    ad_card,
    action_card,
    glossary,
    repair_prompts,
    boundaries,
    one_page_summary,
]


def write_quick_card():
    QUICK_CARD.write_text(
        """# 销售版 Codex 一页速查

## 第一次使用

1. 确认 `amazon-sales-starter` 技能包已安装。
2. 打开一个你平时看的业务页面。
3. 复制下面这句话给 Codex：

```text
只读帮我看当前页面。
按 Amazon 销售入门技能包，先识别页面类型、可见字段、初步判断、风险和下一步建议。
不要直接修改任何内容。
```

## 看结果时检查

- 它有没有说清楚页面类型？
- 它有没有列出看到的字段？
- 它的判断后面有没有证据？
- 它有没有标出风险或不确定项？
- 它有没有默认只读、不直接修改？

## 常用追问

```text
把每个结论后面的证据字段列出来。
```

```text
把支持证据和反证分开写。
```

```text
先生成待确认动作清单，不要执行。
```
""",
        encoding="utf-8",
    )


def build():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    pres = Presentation()
    pres.slide_width = Inches(SLIDE_W)
    pres.slide_height = Inches(SLIDE_H)
    pres.core_properties.title = "销售版 Codex 快速上手卡"
    pres.core_properties.subject = "同事自助阅读版 V4"
    pres.core_properties.author = "Codex"
    for fn in SLIDES:
        fn(pres)
    pres.save(PPTX_PATH)
    write_quick_card()
    print(PPTX_PATH)
    print(QUICK_CARD)


if __name__ == "__main__":
    build()
