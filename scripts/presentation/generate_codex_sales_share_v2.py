from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


OUT_DIR = Path(r"D:\ad-ops-workbench\outputs\codex_sales_share_v2")
PPTX_PATH = OUT_DIR / "sales_codex_share_v2_ready_to_present.pptx"
NOTES_PATH = OUT_DIR / "sales_codex_share_v2_speaker_script.md"

SLIDE_W = 13.333
SLIDE_H = 7.5
FONT_CN = "Microsoft YaHei"
FONT_LATIN = "Arial"

COLORS = {
    "bg": "F6F8FB",
    "paper": "FFFFFF",
    "ink": "0F172A",
    "muted": "64748B",
    "subtle": "94A3B8",
    "line": "D7DEE9",
    "navy": "111827",
    "blue": "2563EB",
    "blue_dark": "1D4ED8",
    "blue_soft": "EAF1FF",
    "teal": "0F766E",
    "teal_soft": "E6F7F4",
    "orange": "F97316",
    "orange_soft": "FFF2E8",
    "green": "059669",
    "green_soft": "E9FBF4",
    "purple": "7C3AED",
    "purple_soft": "F3EAFE",
    "red": "DC2626",
    "red_soft": "FEECEC",
    "yellow": "F59E0B",
    "yellow_soft": "FFF7D6",
}


def rgb(value):
    return RGBColor.from_string(value.replace("#", ""))


def set_font(run, size=16, color="0F172A", bold=False):
    run.font.name = FONT_CN
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    rpr = run._r.get_or_add_rPr()
    rfonts = rpr.find(qn("a:rFonts"))
    if rfonts is None:
        rfonts = OxmlElement("a:rFonts")
        rpr.append(rfonts)
    rfonts.set(qn("a:latin"), FONT_LATIN)
    rfonts.set(qn("a:ea"), FONT_CN)


def clear_frame(tf):
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)


def text_box(slide, text, x, y, w, h, size=16, color="0F172A", bold=False,
             align="left", valign="top", margin=0.02):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    clear_frame(tf)
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    if valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    run = p.add_run()
    run.text = text
    set_font(run, size=size, color=color, bold=bold)
    return shape


def rect(slide, x, y, w, h, fill="FFFFFF", line=None, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    if line:
        shp.line.color.rgb = rgb(line)
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    return shp


def line(slide, x1, y1, x2, y2, color="D7DEE9", width=1.0):
    shp = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    shp.line.color.rgb = rgb(color)
    shp.line.width = Pt(width)
    return shp


def set_shape_text(shape, title=None, body=None, title_size=16, body_size=12,
                   title_color="0F172A", body_color="64748B", title_bold=True,
                   align="left", valign="top", margin=0.16):
    tf = shape.text_frame
    clear_frame(tf)
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    if valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    first = True
    if title:
        p = tf.paragraphs[0]
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
        p.space_after = Pt(7 if body else 0)
        r = p.add_run()
        r.text = title
        set_font(r, size=title_size, color=title_color, bold=title_bold)
        first = False
    if body:
        for i, part in enumerate(body.split("\n")):
            p = tf.paragraphs[0] if first and i == 0 else tf.add_paragraph()
            p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
            p.space_after = Pt(4)
            r = p.add_run()
            r.text = part
            set_font(r, size=body_size, color=body_color, bold=False)
        first = False
    return shape


def pill(slide, text, x, y, w, fill="EAF1FF", color="1D4ED8", size=10.5):
    shp = rect(slide, x, y, w, 0.36, fill=fill, line=None, radius=True)
    set_shape_text(shp, title=text, title_size=size, title_color=color, title_bold=True, align="center", valign="middle", margin=0.02)
    return shp


def card(slide, x, y, w, h, title, body, accent="2563EB", fill="FFFFFF"):
    base = rect(slide, x, y, w, h, fill=fill, line="D7DEE9", radius=True)
    rect(slide, x, y, 0.08, h, fill=accent, line=None, radius=False)
    set_shape_text(base, title=title, body=body, title_size=15, body_size=11.8, title_color="0F172A", body_color="475569", margin=0.22)
    return base


def small_card(slide, x, y, w, h, title, body, fill, accent):
    base = rect(slide, x, y, w, h, fill=fill, line=None, radius=True)
    set_shape_text(base, title=title, body=body, title_size=13, body_size=10.5, title_color=accent, body_color="475569", margin=0.14)
    return base


def background(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(COLORS["bg"])


def header(slide, page, title, subtitle=None, tag=None):
    if tag:
        pill(slide, tag, 0.68, 0.42, 1.28, fill=COLORS["blue_soft"], color=COLORS["blue_dark"], size=9.5)
    text_box(slide, title, 0.68, 0.82, 11.45, 0.42, size=24, bold=True)
    if subtitle:
        text_box(slide, subtitle, 0.7, 1.28, 11.5, 0.28, size=11.5, color=COLORS["muted"])
    footer(slide, page)


def footer(slide, page):
    line(slide, 0.68, 7.06, 12.65, 7.06, color="E2E8F0", width=0.6)
    text_box(slide, "销售 Codex 分享", 0.68, 7.16, 2.8, 0.15, size=8, color=COLORS["subtle"])
    text_box(slide, f"{page:02d}", 12.2, 7.16, 0.45, 0.15, size=8, color=COLORS["subtle"], align="right")


def browser_mock(slide, x, y, w, h, title, rows, active_idx=None):
    base = rect(slide, x, y, w, h, fill="FFFFFF", line="D7DEE9", radius=True)
    top = rect(slide, x, y, w, 0.48, fill="EEF2F7", line=None, radius=True)
    for idx, c in enumerate(["EF4444", "F59E0B", "10B981"]):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.22 + idx * 0.22), Inches(y + 0.17), Inches(0.09), Inches(0.09))
        dot.fill.solid()
        dot.fill.fore_color.rgb = rgb(c)
        dot.line.fill.background()
    text_box(slide, title, x + 0.88, y + 0.13, w - 1.15, 0.19, size=8.8, color="475569", valign="middle")
    for i, row in enumerate(rows):
        fy = y + 0.75 + i * 0.52
        fill = "EAF1FF" if active_idx == i else "FFFFFF"
        r = rect(slide, x + 0.28, fy, w - 0.56, 0.34, fill=fill, line="E2E8F0", radius=True)
        set_shape_text(r, title=row, title_size=9.3, title_color="334155", title_bold=False, align="left", valign="middle", margin=0.08)
    return base


def step_circle(slide, x, y, num, color):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.46), Inches(0.46))
    dot.fill.solid()
    dot.fill.fore_color.rgb = rgb(color)
    dot.line.fill.background()
    text_box(slide, str(num), x, y + 0.08, 0.46, 0.14, size=10.5, color="FFFFFF", bold=True, align="center")
    return dot


def cover(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    rect(slide, 0, 0, 5.2, SLIDE_H, fill="111827", line=None, radius=False)
    text_box(slide, "别从头做一遍", 0.68, 1.18, 4.0, 0.55, size=30, color="FFFFFF", bold=True)
    text_box(slide, "让 Codex 先跑第一遍", 0.68, 1.88, 4.25, 0.48, size=24, color="93C5FD", bold=True)
    text_box(slide, "销售同事的 Codex 入门分享", 0.72, 3.0, 3.9, 0.26, size=12.5, color="CBD5E1")
    text_box(slide, "今天只讲一件事：遇到重复步骤，怎么让它先帮你看、找、整理、列清单。", 0.72, 3.48, 3.7, 0.68, size=14, color="E5E7EB")
    pill(slide, "现场可跟着试", 0.72, 5.2, 1.58, fill="EAF1FF", color="1D4ED8")
    pill(slide, "先只读", 2.45, 5.2, 1.0, fill="E6F7F4", color="0F766E")

    browser_mock(slide, 6.08, 0.86, 5.95, 2.95, "已登录 Chrome 页面", ["广告系统 / SIF / SellerInventory", "SKU、ASIN、搜索词、竞品页面", "页面可见数据", "待确认动作清单"], active_idx=2)
    labels = [("你说一句", "2563EB"), ("它去看", "F97316"), ("给清单", "059669")]
    for i, (txt, color) in enumerate(labels):
        box = rect(slide, 6.36 + i * 1.76, 4.62, 1.28, 0.72, fill="FFFFFF", line="D7DEE9", radius=True)
        set_shape_text(box, title=txt, title_size=13.5, title_color=color, title_bold=True, align="center", valign="middle", margin=0.04)
        if i < 2:
            line(slide, 7.72 + i * 1.76, 4.98, 8.04 + i * 1.76, 4.98, color="CBD5E1", width=1.4)
    text_box(slide, "不是换系统，也不是多学一套流程。", 6.32, 5.78, 5.45, 0.25, size=13, color=COLORS["muted"], align="center")
    footer(slide, 1)


def promise(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    header(slide, 2, "今天不讲高深 AI，只解决一个销售日常问题", "很多事你都会做，但每次都要重新打开页面、找证据、整理判断。", "开场")
    left = rect(slide, 0.92, 2.0, 5.35, 3.65, fill="FFFFFF", line="D7DEE9", radius=True)
    set_shape_text(left, "以前", "打开页面\n自己找字段\n自己比较和判断\n再整理成文字", title_size=18, body_size=17, title_color=COLORS["red"], margin=0.32)
    right = rect(slide, 7.05, 2.0, 5.35, 3.65, fill="FFFFFF", line="D7DEE9", radius=True)
    set_shape_text(right, "现在先试", "把重复步骤说给 Codex\n让它先看一遍\n输出证据和清单\n你再确认", title_size=18, body_size=17, title_color=COLORS["green"], margin=0.32)
    line(slide, 6.42, 3.78, 6.86, 3.78, color=COLORS["orange"], width=2.2)
    text_box(slide, "目标：少掉第一遍重复劳动", 4.8, 5.95, 3.9, 0.32, size=14.5, color=COLORS["orange"], bold=True, align="center")


def difference(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    header(slide, 3, "它不是又一个聊天框", "只让 AI 改文案，ChatGPT 也可以。Codex 的价值在于：它可以跟着你的工作页面走。", "认知")
    card(slide, 0.9, 2.0, 5.35, 3.8, "ChatGPT：你喂材料", "你复制资料、粘贴表格、描述背景。\n它帮你总结、润色、写汇报。", accent="64748B")
    card(slide, 7.05, 2.0, 5.35, 3.8, "Codex：它先去看", "在你的授权和权限范围内，读取页面和文件。\n它帮你找证据、跑步骤、列清单。", accent=COLORS["blue"])
    text_box(slide, "一句话记住：ChatGPT 更会写，Codex 更会跟着流程做。", 2.0, 6.12, 9.3, 0.34, size=15, color=COLORS["ink"], bold=True, align="center")


def chrome_ability(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    header(slide, 4, "最大变化：它能看你已经登录的 Chrome 页面", "不是绕权限，不是后台导出。它看的就是你本来能看的页面。", "能力")
    browser_mock(slide, 0.92, 1.9, 5.65, 3.7, "你的 Chrome 工作页面", ["广告系统：SKU 广告表现", "SIF：关键词和市场证据", "SellerInventory：库存、利润、状态", "Amazon 前台：竞品和页面表达"], active_idx=0)
    items = [
        ("只能看本人权限", "看不到你没有权限的东西", COLORS["blue_soft"], COLORS["blue"]),
        ("先按只读使用", "新手先不让它修改", COLORS["teal_soft"], COLORS["teal"]),
        ("结果可以查验", "每个判断要能回到页面核对", COLORS["orange_soft"], COLORS["orange"]),
        ("成熟后再授权动作", "先清单，后确认，再执行", COLORS["purple_soft"], COLORS["purple"]),
    ]
    for i, (t, b, fill, accent) in enumerate(items):
        small_card(slide, 7.05, 1.9 + i * 0.9, 5.1, 0.7, t, b, fill, accent)
    text_box(slide, "这就是为什么它不只是“帮你写”，而是可以帮你少跑一段流程。", 1.2, 6.12, 10.9, 0.34, size=14.5, color=COLORS["ink"], bold=True, align="center")


def when_to_use(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    header(slide, 5, "什么时候该想到 Codex？看这四个信号", "不要背复杂场景。只要一个任务符合这些信号，就可以让它先跑第一遍。", "判断")
    signals = [
        ("每周都做", "周会前、每个 SKU、每次促销前都要看"),
        ("步骤相似", "先看页面，再找字段，再比较，再写结论"),
        ("有固定方法", "你心里知道先看什么、什么情况算异常"),
        ("结果能核对", "它给的证据可以回到页面或文件里确认"),
    ]
    colors = [(COLORS["blue_soft"], COLORS["blue"]), (COLORS["teal_soft"], COLORS["teal"]), (COLORS["orange_soft"], COLORS["orange"]), (COLORS["purple_soft"], COLORS["purple"])]
    for i, (title, body) in enumerate(signals):
        x = 0.92 + (i % 2) * 6.05
        y = 2.05 + (i // 2) * 1.55
        fill, accent = colors[i]
        box = rect(slide, x, y, 5.25, 1.15, fill="FFFFFF", line="D7DEE9", radius=True)
        step_circle(slide, x + 0.24, y + 0.34, i + 1, accent)
        text_box(slide, title, x + 0.9, y + 0.25, 2.1, 0.25, size=15, color=accent, bold=True)
        text_box(slide, body, x + 0.9, y + 0.66, 3.9, 0.28, size=11.5, color=COLORS["muted"])
    rect(slide, 2.1, 5.8, 9.15, 0.58, fill="111827", line=None, radius=True)
    text_box(slide, "你不用学 AI，你只要把自己的重复方法说清楚。", 2.45, 5.95, 8.45, 0.18, size=13, color="FFFFFF", bold=True, align="center")


def scenario_competitor(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    header(slide, 6, "场景一：竞品分析，不看一个链接，要看一片市场", "一个竞品链接不够吸引人。更有价值的是让它一次整理 10 个竞品或一个搜索词下的市场。", "场景")
    blocks = [
        ("你给它", "搜索词\n或 10 个竞品链接\n或 Amazon 前台页面", COLORS["blue"], COLORS["blue_soft"]),
        ("它整理", "价格带\n评论门槛\n主图打法\n卖点结构", COLORS["orange"], COLORS["orange_soft"]),
        ("你拿到", "这个市场靠什么卖\n我们差在哪里\n下一步看什么证据", COLORS["green"], COLORS["green_soft"]),
    ]
    for i, (title, body, accent, fill) in enumerate(blocks):
        x = 0.95 + i * 4.1
        box = rect(slide, x, 2.0, 3.35, 3.25, fill="FFFFFF", line="D7DEE9", radius=True)
        set_shape_text(box, title, body, title_size=17, body_size=15, title_color=accent, body_color="334155", margin=0.28)
        if i < 2:
            line(slide, x + 3.55, 3.58, x + 3.88, 3.58, color="CBD5E1", width=1.5)
    prompt = rect(slide, 1.28, 5.8, 10.78, 0.56, fill="FFFFFF", line="D7DEE9", radius=True)
    set_shape_text(prompt, title="提示词：帮我看这个搜索词下的 10 个竞品，整理价格、评论、主图和卖点差异。", title_size=12.5, title_color=COLORS["ink"], title_bold=True, align="center", valign="middle", margin=0.04)


def scenario_conversion(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    header(slide, 7, "场景二：转化分析，让它先找证据和反证", "你提出判断，它负责把证据先找齐。重点不是让它替你拍脑袋，而是让它少掉你找证据的时间。", "场景")
    quote = rect(slide, 0.95, 2.02, 3.35, 3.5, fill="111827", line=None, radius=True)
    set_shape_text(quote, "人的判断", "我怀疑这个 SKU 不是没流量，而是点击后不转化。", title_size=15, body_size=18, title_color="FFFFFF", body_color="E5E7EB", margin=0.28)
    evidences = [
        ("广告证据", "点击、CTR、花费、订单", COLORS["blue_soft"], COLORS["blue"]),
        ("页面证据", "价格、评分、图片、卖点", COLORS["teal_soft"], COLORS["teal"]),
        ("竞品证据", "同价位对手是否能成交", COLORS["orange_soft"], COLORS["orange"]),
        ("反证", "也可能是流量词不准", COLORS["red_soft"], COLORS["red"]),
    ]
    for i, (title, body, fill, accent) in enumerate(evidences):
        x = 5.0 + (i % 2) * 3.55
        y = 2.05 + (i // 2) * 1.35
        small_card(slide, x, y, 3.05, 0.94, title, body, fill, accent)
    output = rect(slide, 5.0, 5.15, 6.62, 0.78, fill="FFFFFF", line="D7DEE9", radius=True)
    set_shape_text(output, "输出格式", "支持证据 / 反证 / 还缺什么 / 下一步建议", title_size=12, body_size=12.5, title_color=COLORS["blue"], body_color=COLORS["ink"], align="center", valign="middle", margin=0.08)


def method_reuse(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    header(slide, 8, "真正省时间的不是一次分析，而是复用你的方法", "你怎么判断一个 SKU、一个搜索词、一个竞品市场，可以先写成规则，再让 Codex 每次照着跑。", "复用")
    steps = [
        ("写清规则", "先看什么\n异常条件是什么"),
        ("交给 Codex", "按同样顺序\n重复读取和整理"),
        ("检查结果", "证据在哪里\n结论能否核对"),
        ("修改规则", "发现不准\n就改方法论"),
    ]
    accents = [COLORS["blue"], COLORS["orange"], COLORS["green"], COLORS["purple"]]
    for i, (title, body) in enumerate(steps):
        x = 0.95 + i * 3.05
        box = rect(slide, x, 2.18, 2.45, 2.75, fill="FFFFFF", line="D7DEE9", radius=True)
        step_circle(slide, x + 0.24, 2.48, i + 1, accents[i])
        text_box(slide, title, x + 0.24, 3.18, 1.8, 0.25, size=15, color=accents[i], bold=True)
        text_box(slide, body, x + 0.24, 3.7, 1.82, 0.58, size=13.5, color="334155")
        if i < 3:
            line(slide, x + 2.53, 3.55, x + 2.92, 3.55, color="CBD5E1", width=1.3)
    rect(slide, 2.05, 5.75, 9.25, 0.58, fill="111827", line=None, radius=True)
    text_box(slide, "规则透明、结果可查、方法可改，这比“让 AI 猜”靠谱得多。", 2.35, 5.9, 8.65, 0.18, size=13, color="FFFFFF", bold=True, align="center")


def last_mile(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    header(slide, 9, "最后一公里可以做，但必须先有闸门", "不是让它自动乱改。正确顺序是：先只读分析，生成清单，人确认后再执行，并且回读结果。", "动作")
    items = [
        ("先只读", "先让它看页面、找证据、给建议", COLORS["blue_soft"], COLORS["blue"]),
        ("列清单", "动作、对象、理由、风险写清楚", COLORS["orange_soft"], COLORS["orange"]),
        ("人确认", "确认范围、强度、例外项", COLORS["purple_soft"], COLORS["purple"]),
        ("再回读", "改完后重新读取页面，确认落地", COLORS["green_soft"], COLORS["green"]),
    ]
    for i, (title, body, fill, accent) in enumerate(items):
        x = 0.92 + i * 3.07
        box = rect(slide, x, 2.15, 2.45, 2.75, fill="FFFFFF", line="D7DEE9", radius=True)
        set_shape_text(box, title, body, title_size=16, body_size=12.5, title_color=accent, body_color="475569", margin=0.22)
        if i < 3:
            line(slide, x + 2.55, 3.52, x + 2.94, 3.52, color="CBD5E1", width=1.4)
    banner = rect(slide, 1.75, 5.75, 9.85, 0.62, fill="111827", line=None, radius=True)
    set_shape_text(banner, "可以做，但必须有授权、有清单、有回读。", title_size=14, title_color="FFFFFF", title_bold=True, align="center", valign="middle", margin=0.04)


def prompts(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    header(slide, 10, "先不用学复杂提示词，复制这六句话就能开始", "新手最重要的是敢试。每句话都加上“只读”或“不要直接执行”，先把安全边界立住。", "入口")
    data = [
        ("只读诊断", "只读帮我看这个 SKU 的广告页面，先不要做任何修改。"),
        ("竞品市场", "帮我看这个搜索词下的 10 个竞品，整理价格、评论、主图和卖点差异。"),
        ("找证据", "我怀疑这个 SKU 点击后不转化，帮我找支持证据和反证。"),
        ("检查清单", "按我的方法论生成检查清单，结论后面都附证据来源。"),
        ("动作建议", "按规则生成待确认动作清单，不要直接执行。"),
        ("复盘整理", "把这次调整前后的证据整理成复盘草稿，标出还缺什么。"),
    ]
    fills = [COLORS["blue_soft"], COLORS["teal_soft"], COLORS["orange_soft"], COLORS["purple_soft"], COLORS["green_soft"], COLORS["red_soft"]]
    accents = [COLORS["blue"], COLORS["teal"], COLORS["orange"], COLORS["purple"], COLORS["green"], COLORS["red"]]
    for i, (title, body) in enumerate(data):
        x = 0.88 + (i % 2) * 6.1
        y = 1.85 + (i // 2) * 1.45
        small_card(slide, x, y, 5.3, 1.02, title, body, fills[i], accents[i])


def live_practice(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    header(slide, 11, "现场练习：每个人只选一个重复小动作", "不要一开始就做大项目。先选一个你最近真的在重复做的动作，跑一次只读版。", "练习")
    left = rect(slide, 0.95, 1.92, 4.0, 4.15, fill="FFFFFF", line="D7DEE9", radius=True)
    set_shape_text(left, "选一个任务", "我最近总是在重复看什么？\n\n竞品？\n广告？\n转化？\n库存状态？", title_size=17, body_size=14.5, title_color=COLORS["blue"], body_color="334155", margin=0.28)
    middle = rect(slide, 5.22, 1.92, 3.1, 4.15, fill="FFFFFF", line="D7DEE9", radius=True)
    set_shape_text(middle, "加一句边界", "只读。\n先不要修改。\n结论后面附证据。", title_size=17, body_size=15.5, title_color=COLORS["orange"], body_color="334155", margin=0.28)
    right = rect(slide, 8.6, 1.92, 3.78, 4.15, fill="111827", line=None, radius=True)
    set_shape_text(right, "看结果", "它是不是少帮你跑了一遍？\n证据能不能查？\n下次能不能复用？", title_size=17, body_size=14.5, title_color="FFFFFF", body_color="E5E7EB", margin=0.28)
    text_box(slide, "这一步的成功标准不是“AI 很聪明”，而是“我愿不愿意下次再叫它一次”。", 1.55, 6.35, 10.3, 0.28, size=13.5, color=COLORS["ink"], bold=True, align="center")


def boundaries(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    header(slide, 12, "放心试，但要记住四条边界", "边界讲清楚，大家才敢用，组内也更容易推广。", "安全")
    items = [
        ("不越权", "只看你本来有权限看的内容", COLORS["blue_soft"], COLORS["blue"]),
        ("不乱改", "新手默认只读，动作要确认", COLORS["red_soft"], COLORS["red"]),
        ("不迷信", "结论必须能回到证据核对", COLORS["orange_soft"], COLORS["orange"]),
        ("不外传", "敏感资料按公司要求处理", COLORS["teal_soft"], COLORS["teal"]),
    ]
    for i, (title, body, fill, accent) in enumerate(items):
        x = 0.95 + (i % 2) * 6.08
        y = 2.0 + (i // 2) * 1.55
        box = rect(slide, x, y, 5.35, 1.15, fill=fill, line=None, radius=True)
        set_shape_text(box, title, body, title_size=15, body_size=11.5, title_color=accent, body_color="334155", margin=0.18)
    rect(slide, 2.2, 5.7, 8.95, 0.68, fill="FFFFFF", line="D7DEE9", radius=True)
    text_box(slide, "一句话：让它帮你先跑流程，不是让它替你担责任。", 2.5, 5.9, 8.35, 0.2, size=13.5, color=COLORS["ink"], bold=True, align="center")


def install(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    header(slide, 13, "想试的人，会后按公司指南准备", "今天先看效果。真正开始时，不用自己研究安装，按公司指南走，卡住找协助。", "会后")
    steps = [
        ("1", "申请公司版 Codex", "按公司流程申请使用权限。"),
        ("2", "打开企业微信指南", "文档名：Codex安装&使用教程。"),
        ("3", "按提示完成安装", "复制安装 key，按提示执行，重启软件。"),
        ("4", "登录常用 Chrome 页面", "先从只读任务开始试。"),
    ]
    for i, (num, title, body) in enumerate(steps):
        x = 0.95 + (i % 2) * 6.08
        y = 1.95 + (i // 2) * 1.35
        box = rect(slide, x, y, 5.35, 0.95, fill="FFFFFF", line="D7DEE9", radius=True)
        step_circle(slide, x + 0.22, y + 0.25, num, [COLORS["blue"], COLORS["orange"], COLORS["green"], COLORS["purple"]][i])
        text_box(slide, title, x + 0.88, y + 0.18, 3.8, 0.22, size=14.2, color=COLORS["ink"], bold=True)
        text_box(slide, body, x + 0.88, y + 0.52, 4.1, 0.2, size=10.8, color=COLORS["muted"])
    assist = rect(slide, 2.0, 5.35, 9.35, 0.72, fill="111827", line=None, radius=True)
    set_shape_text(assist, "卡住不用硬扛：联系方燕 / 张文琦远程协助。", title_size=14, title_color="FFFFFF", title_bold=True, align="center", valign="middle", margin=0.04)
    text_box(slide, "提示：不要在开场就讲安装。先让大家看到效果，最后再告诉想试的人怎么开始。", 2.1, 6.35, 9.1, 0.24, size=10.5, color=COLORS["muted"], align="center")


def closing(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    rect(slide, 0.8, 0.78, 11.75, 5.55, fill="111827", line=None, radius=True)
    text_box(slide, "今天先不用变成高手", 1.35, 1.55, 10.7, 0.55, size=30, color="FFFFFF", bold=True, align="center")
    text_box(slide, "找一个最近重复做的事，让 Codex 先帮你做一小步。", 1.75, 2.36, 9.9, 0.36, size=18, color="E5E7EB", align="center")
    labels = [("看页面", COLORS["blue_soft"], COLORS["blue"]), ("找证据", COLORS["orange_soft"], COLORS["orange"]), ("列清单", COLORS["green_soft"], COLORS["green"])]
    for i, (txt, fill, accent) in enumerate(labels):
        box = rect(slide, 2.05 + i * 3.2, 3.55, 2.35, 1.0, fill="FFFFFF", line=None, radius=True)
        set_shape_text(box, txt, None, title_size=17, title_color=accent, title_bold=True, align="center", valign="middle", margin=0.04)
    text_box(slide, "重复动作越多，越值得让工具先跑第一遍。", 1.55, 5.25, 10.25, 0.28, size=13.5, color="CBD5E1", align="center")
    footer(slide, 14)


SLIDES = [
    cover,
    promise,
    difference,
    chrome_ability,
    when_to_use,
    scenario_competitor,
    scenario_conversion,
    method_reuse,
    last_mile,
    prompts,
    live_practice,
    boundaries,
    install,
    closing,
]


NOTES = [
    ("封面", "今天只讲一个核心：遇到重复步骤，不要每次从头做，让 Codex 先跑第一遍。"),
    ("开场痛点", "承认大家不是不会做，而是重复动作太多。先建立共鸣。"),
    ("区别", "把 ChatGPT 和 Codex 的区别讲在工作流上，不讲模型参数。"),
    ("Chrome 能力", "强调本人权限、已登录页面、先只读。不要提抓包或后台路径。"),
    ("何时使用", "让大家记住四个信号：每周做、步骤相似、有方法、可核对。"),
    ("竞品场景", "重点是看一片市场，不是单个链接。讲完直接读提示词。"),
    ("转化场景", "人提出判断，Codex 找支持证据和反证。"),
    ("方法复用", "引导大家把自己的判断方法写成规则，后续可以复用。"),
    ("最后一公里", "可以动作，但必须先有授权、有清单、有回读。"),
    ("六句提示词", "这页可以让大家拍照或会后复制，降低试用门槛。"),
    ("现场练习", "让每个人选一个最近重复动作，强调只读边界。"),
    ("安全边界", "边界越清楚，大家越敢试。"),
    ("会后安装", "安装放最后，卡住找协助，不要让安装变成心理负担。"),
    ("收尾", "只要求试一个小动作，不要求大家当天变成高手。"),
]


def build():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    pres = Presentation()
    pres.slide_width = Inches(SLIDE_W)
    pres.slide_height = Inches(SLIDE_H)
    pres.core_properties.title = "销售怎么开始用 Codex"
    pres.core_properties.subject = "正式分享课件 V2"
    pres.core_properties.author = "Codex"
    for fn in SLIDES:
        fn(pres)
    pres.save(PPTX_PATH)

    lines = ["# 销售 Codex 分享 V2 讲稿要点", ""]
    for idx, (title, body) in enumerate(NOTES, start=1):
        lines.append(f"## Slide {idx:02d} - {title}")
        lines.append(body)
        lines.append("")
    NOTES_PATH.write_text("\n".join(lines), encoding="utf-8")
    print(PPTX_PATH)
    print(NOTES_PATH)


if __name__ == "__main__":
    build()
