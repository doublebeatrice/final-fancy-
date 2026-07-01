from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from generate_codex_sales_share_v2 import (
    COLORS,
    SLIDE_H,
    SLIDE_W,
    background,
    browser_mock,
    card,
    footer,
    header,
    line,
    rect,
    set_shape_text,
    small_card,
    step_circle,
    text_box,
)


OUT_DIR = Path(r"D:\ad-ops-workbench\outputs\codex_sales_share_v5")
PPTX_PATH = OUT_DIR / "sales_codex_intro_and_start_v5_for_colleagues.pptx"
QUICK_CARD = OUT_DIR / "sales_codex_colleague_quick_reference.md"


def title(slide, page, main, sub=None, tag=None):
    header(slide, page, main, sub, tag)


def dark_prompt(slide, x, y, w, h, title_text, body_text):
    box = rect(slide, x, y, w, h, fill="111827", line=None, radius=True)
    set_shape_text(
        box,
        title_text,
        body_text,
        title_size=15.5,
        body_size=13.2,
        title_color="FFFFFF",
        body_color="E5E7EB",
        margin=0.22,
    )
    return box


def cover(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    rect(slide, 0.82, 0.78, 11.7, 5.6, fill="111827", line=None, radius=True)
    text_box(slide, "销售同事第一次认识 Codex", 1.35, 1.35, 10.65, 0.55, size=31, color="FFFFFF", bold=True, align="center")
    text_box(slide, "它是什么，为什么有用，第一天怎么开始", 1.35, 2.13, 10.65, 0.42, size=20, color="93C5FD", bold=True, align="center")
    text_box(slide, "这份材料默认你以前没用过 Codex，也不知道它和 ChatGPT 有什么区别。", 2.05, 3.02, 9.25, 0.35, size=15, color="E5E7EB", align="center")
    boxes = [("先认识", COLORS["blue"]), ("再装技能包", COLORS["orange"]), ("最后只读试一次", COLORS["green"])]
    for i, (txt, color) in enumerate(boxes):
        b = rect(slide, 2.0 + i * 3.2, 4.35, 2.35, 0.85, fill="FFFFFF", line=None, radius=True)
        set_shape_text(b, txt, None, title_size=15.5, title_color=color, title_bold=True, align="center", valign="middle", margin=0.04)
    footer(slide, 1)


def what_is_codex(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    title(slide, 2, "Codex 是什么？先用一句话理解", "它不是一个新业务系统，也不是只会聊天的软件。", "认识")
    left = rect(slide, 0.95, 2.0, 5.35, 3.65, fill="FFFFFF", line="D7DEE9", radius=True)
    set_shape_text(left, "一句话", "Codex 是一个能在你授权范围内，帮你看页面、读文件、整理证据、生成清单的 AI 工作助手。", title_size=18, body_size=16, title_color=COLORS["blue"], body_color="334155", margin=0.3)
    right = rect(slide, 7.02, 2.0, 5.35, 3.65, fill="111827", line=None, radius=True)
    set_shape_text(right, "销售可以先这样理解", "你原来要自己打开页面、找字段、整理判断。\n\n现在可以先让它只读跑第一遍。", title_size=18, body_size=16, title_color="FFFFFF", body_color="E5E7EB", margin=0.3)
    text_box(slide, "重点：不是让它替你负责，而是让它先帮你少做一遍重复动作。", 1.8, 6.12, 9.8, 0.28, size=14.5, color=COLORS["ink"], bold=True, align="center")


def chatgpt_difference(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    title(slide, 3, "它和 ChatGPT 有什么不一样？", "如果只是写汇报、润色文字，ChatGPT 也可以。Codex 的特点是能跟你的工作流更近。", "认识")
    card(slide, 0.95, 2.0, 5.35, 3.65, "ChatGPT 更像聊天和写作助手", "你把材料复制进去。\n它帮你总结、改写、润色、写复盘。", accent="64748B")
    card(slide, 7.05, 2.0, 5.35, 3.65, "Codex 更像工作流助手", "它可以看你打开的页面和文件。\n它帮你找证据、跑步骤、列清单。", accent=COLORS["blue"])
    text_box(slide, "销售最先用它，不是为了写得更漂亮，而是为了少重复打开、查找、整理。", 1.55, 6.12, 10.15, 0.28, size=14, color=COLORS["ink"], bold=True, align="center")


def why_sales_care(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    title(slide, 4, "为什么销售同事要关心它？", "因为销售工作里有很多“每次都差不多，但每次都要重新做”的步骤。", "场景感")
    items = [
        ("看产品", "价格、评价、图片、卖点、状态"),
        ("看广告", "曝光、点击、花费、订单、ACOS、CVR"),
        ("看库存利润", "FBA、库存天数、滞销风险、空海利润"),
        ("看竞品市场", "价格带、评论门槛、主图打法、卖点差异"),
    ]
    colors = [COLORS["blue"], COLORS["orange"], COLORS["green"], COLORS["purple"]]
    fills = [COLORS["blue_soft"], COLORS["orange_soft"], COLORS["green_soft"], COLORS["purple_soft"]]
    for i, (t, b) in enumerate(items):
        x = 0.95 + (i % 2) * 6.08
        y = 1.95 + (i // 2) * 1.55
        small_card(slide, x, y, 5.35, 1.08, t, b, fills[i], colors[i])
    rect(slide, 2.1, 5.75, 9.15, 0.62, fill="111827", line=None, radius=True)
    text_box(slide, "只要一件事你经常重复做，就可以先让 Codex 跑第一遍。", 2.4, 5.93, 8.55, 0.16, size=13.2, color="FFFFFF", bold=True, align="center")


def blank_codex(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    title(slide, 5, "但要先说清楚：刚装好的 Codex 是空白的", "它不会天然懂 Amazon 销售，也不会天然知道公司内部业务词。", "真实情况")
    left = rect(slide, 0.95, 2.0, 5.35, 3.6, fill="FFFFFF", line="D7DEE9", radius=True)
    set_shape_text(left, "它不知道", "SKU / ASIN 是什么关系\nACOS、CVR 怎么看\n滞销风险意味着什么\n空运和海运利润怎么理解\n哪些动作不能直接做", title_size=17, body_size=14.5, title_color=COLORS["red"], body_color="334155", margin=0.3)
    right = rect(slide, 7.05, 2.0, 5.35, 3.6, fill="FFFFFF", line="D7DEE9", radius=True)
    set_shape_text(right, "所以不要这样用", "不要一开始就问：\n“帮我分析这个 SKU。”\n\n它可能会猜，也可能问你一堆基础问题。", title_size=17, body_size=14.5, title_color=COLORS["orange"], body_color="334155", margin=0.3)
    text_box(slide, "解决办法：先给它装销售技能包。", 3.25, 6.05, 6.8, 0.3, size=15, color=COLORS["blue"], bold=True, align="center")


def sales_skill_pack(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    title(slide, 6, "销售技能包是做什么的？", "它相当于先给 Codex 一份销售通用手册，避免每个人从零解释。", "技能包")
    center = rect(slide, 5.2, 2.0, 2.9, 3.4, fill="111827", line=None, radius=True)
    set_shape_text(center, "销售\n技能包", "amazon-sales-starter", title_size=26, body_size=12, title_color="FFFFFF", body_color="CBD5E1", align="center", valign="middle", margin=0.05)
    left_items = [("常用职责", "产品 / 广告 / 库存 / 竞品"), ("基础术语", "SKU / ASIN / ACOS / CVR / FBA")]
    right_items = [("判断顺序", "先目标，再产品，再流量"), ("安全边界", "默认只读，不直接修改")]
    for i, (t, b) in enumerate(left_items):
        small_card(slide, 0.95, 2.15 + i * 1.45, 3.75, 1.0, t, b, COLORS["blue_soft"], COLORS["blue"])
        line(slide, 4.82, 2.65 + i * 1.45, 5.12, 2.65 + i * 1.45, color="CBD5E1", width=1.4)
    for i, (t, b) in enumerate(right_items):
        small_card(slide, 8.55, 2.15 + i * 1.45, 3.75, 1.0, t, b, COLORS["green_soft"], COLORS["green"])
        line(slide, 8.18, 2.65 + i * 1.45, 8.48, 2.65 + i * 1.45, color="CBD5E1", width=1.4)
    text_box(slide, "装完技能包后，Codex 不是懂你所有业务细节，但至少不会从零开始。", 1.7, 6.15, 9.95, 0.26, size=13.5, color=COLORS["ink"], bold=True, align="center")


def after_skill(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    title(slide, 7, "装好技能包后，你可以期待它先做到这些", "先别追求全自动。第一阶段只看它能不能把页面和证据整理清楚。", "预期")
    items = [
        ("识别页面", "这是广告页、库存页、Amazon 前台还是竞品页面"),
        ("解释字段", "常见广告、库存、产品字段先按销售语境理解"),
        ("归类问题", "曝光不足、点击不足、转化问题、效率问题"),
        ("列证据", "每个判断后面对应可见字段或页面信息"),
        ("提醒边界", "不确定项标出来，动作先生成待确认清单"),
    ]
    for i, (t, b) in enumerate(items):
        x = 0.95 + (i % 2) * 6.05
        y = 1.72 + (i // 2) * 1.22
        w = 5.25 if i < 4 else 11.28
        if i == 4:
            x = 0.95
        small_card(slide, x, y, w, 0.82, t, b, [COLORS["blue_soft"], COLORS["teal_soft"], COLORS["orange_soft"], COLORS["purple_soft"], COLORS["green_soft"]][i], [COLORS["blue"], COLORS["teal"], COLORS["orange"], COLORS["purple"], COLORS["green"]][i])


def install(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    title(slide, 8, "准备工作：先安装 Codex，再放入销售技能包", "如果你还没安装 Codex，先按公司企业微信指南完成安装。技能包可以找助理协助放好。", "准备")
    steps = [
        ("1", "安装公司版 Codex", "按企业微信《Codex安装&使用教程》。"),
        ("2", "拿到技能包文件夹", "文件夹名：amazon-sales-starter。"),
        ("3", "放进 skills 目录", "常见路径：C:\\Users\\<用户名>\\.codex\\skills\\"),
        ("4", "重启 Codex", "然后打开一个你平时看的业务页面。"),
    ]
    for i, (num, t, b) in enumerate(steps):
        x = 0.95 + (i % 2) * 6.08
        y = 1.95 + (i // 2) * 1.35
        box = rect(slide, x, y, 5.35, 0.95, fill="FFFFFF", line="D7DEE9", radius=True)
        step_circle(slide, x + 0.22, y + 0.25, num, [COLORS["blue"], COLORS["orange"], COLORS["green"], COLORS["purple"]][i])
        text_box(slide, t, x + 0.88, y + 0.18, 3.8, 0.22, size=14.2, color=COLORS["ink"], bold=True)
        text_box(slide, b, x + 0.88, y + 0.52, 4.1, 0.2, size=10.8, color=COLORS["muted"])
    rect(slide, 2.0, 5.65, 9.35, 0.72, fill="111827", line=None, radius=True)
    text_box(slide, "找不到目录或安装卡住，不要硬研究，找助理远程协助。", 2.25, 5.88, 8.85, 0.16, size=13.5, color="FFFFFF", bold=True, align="center")


def first_three_steps(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    title(slide, 9, "第一次使用，只做三步", "不要一上来做复杂任务。先在一个你熟悉的页面上试一次只读。", "开始")
    steps = [
        ("1", "打开页面", "广告页、库存页、Amazon 前台、竞品页面都可以"),
        ("2", "复制提示词", "先只读，不直接修改"),
        ("3", "检查结果", "看它是否列出证据、风险和下一步"),
    ]
    for i, (num, t, b) in enumerate(steps):
        x = 1.0 + i * 4.05
        box = rect(slide, x, 2.0, 3.35, 3.45, fill="FFFFFF", line="D7DEE9", radius=True)
        step_circle(slide, x + 0.3, 2.38, num, [COLORS["blue"], COLORS["orange"], COLORS["green"]][i])
        text_box(slide, t, x + 0.3, 3.15, 2.5, 0.3, size=17, color=COLORS["ink"], bold=True)
        text_box(slide, b, x + 0.3, 3.78, 2.55, 0.6, size=13.5, color=COLORS["muted"])
        if i < 2:
            line(slide, x + 3.52, 3.68, x + 3.84, 3.68, color="CBD5E1", width=1.5)


def first_prompt(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    title(slide, 10, "第一句测试：直接复制这段", "打开任意一个你平时看的业务页面，然后把下面这段发给 Codex。", "开始")
    dark_prompt(
        slide,
        1.2,
        1.8,
        10.9,
        2.75,
        "复制给 Codex",
        "只读帮我看当前页面。\n按 Amazon 销售入门技能包，先识别页面类型、可见字段、初步判断、风险和下一步建议。\n不要直接修改任何内容。",
    )
    rows = [
        ("它应该先说", "这是什么页面、看到了哪些对象和字段"),
        ("它不应该做", "不应该直接改广告、改价格、改 listing"),
        ("你要检查", "结论后面有没有证据，是否有不确定项"),
    ]
    for i, (t, b) in enumerate(rows):
        small_card(slide, 1.2 + i * 3.65, 5.1, 3.1, 0.85, t, b, [COLORS["blue_soft"], COLORS["red_soft"], COLORS["green_soft"]][i], [COLORS["blue"], COLORS["red"], COLORS["green"]][i])


def expected_output(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    title(slide, 11, "第一次输出，重点看这五块", "不要只看它说得顺不顺，要看它有没有证据、风险和下一步。", "检查")
    sections = [
        ("看到的信息", "页面类型、SKU/ASIN/搜索词、可见字段"),
        ("初步判断", "问题更像曝光、点击、转化还是效率"),
        ("支持证据", "每个判断对应哪个字段或页面信息"),
        ("风险/不确定", "哪些字段看不懂、哪些结论样本不足"),
        ("下一步建议", "先补证据，还是生成待确认动作清单"),
    ]
    for i, (t, b) in enumerate(sections):
        x = 0.95 + (i % 2) * 6.05
        y = 1.72 + (i // 2) * 1.22
        w = 5.25
        if i == 4:
            x = 0.95
            w = 11.28
        small_card(slide, x, y, w, 0.82, t, b, [COLORS["blue_soft"], COLORS["teal_soft"], COLORS["orange_soft"], COLORS["red_soft"], COLORS["green_soft"]][i], [COLORS["blue"], COLORS["teal"], COLORS["orange"], COLORS["red"], COLORS["green"]][i])
    rect(slide, 2.15, 6.0, 9.05, 0.55, fill="111827", line=None, radius=True)
    text_box(slide, "如果它只给一句“建议优化”，就让它补证据。", 2.45, 6.15, 8.45, 0.16, size=12.5, color="FFFFFF", bold=True, align="center")


def current_page(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    title(slide, 12, "最推荐从“当前页面只读”开始", "不需要导表，不需要截图很多张。先让它看你已经打开的页面。", "用法")
    browser_mock(slide, 0.95, 1.82, 5.6, 3.75, "已打开 Chrome 页面", ["广告页面", "库存/产品页面", "Amazon 前台页面", "竞品搜索结果"], active_idx=1)
    right = rect(slide, 7.05, 1.82, 5.2, 3.75, fill="FFFFFF", line="D7DEE9", radius=True)
    set_shape_text(right, "你要做", "打开页面\n确认自己能看到数据\n复制只读提示词\n检查它输出的证据", title_size=18, body_size=16, title_color=COLORS["green"], body_color="334155", margin=0.3)
    text_box(slide, "第一阶段只追求：它帮你少看一遍、少整理一遍。", 2.0, 6.12, 9.35, 0.28, size=14, color=COLORS["ink"], bold=True, align="center")


def competitor(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    title(slide, 13, "用法 1：竞品市场快速整理", "适合：你要看一个搜索词、一组竞品、一个新品机会。", "用法")
    dark_prompt(slide, 0.95, 1.85, 5.35, 2.1, "复制提示词", "帮我看这个搜索词下的 10 个竞品，整理价格、评论、主图和卖点差异。最后告诉我我们能不能打、还缺什么证据。")
    right = rect(slide, 7.0, 1.85, 5.35, 2.1, fill="FFFFFF", line="D7DEE9", radius=True)
    set_shape_text(right, "你会拿到", "价格带\n评论门槛\n主图打法\n卖点结构\n差距和机会点", title_size=16, body_size=13.8, title_color=COLORS["blue"], body_color="334155", margin=0.22)
    checks = [("检查 1", "是否只看了单个链接？"), ("检查 2", "是否有“我们能不能打”的结论？"), ("检查 3", "是否标出还缺什么证据？")]
    for i, (t, b) in enumerate(checks):
        small_card(slide, 0.95 + i * 4.05, 4.8, 3.35, 0.9, t, b, COLORS["blue_soft"], COLORS["blue"])


def conversion(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    title(slide, 14, "用法 2：点击后不转化，先找证据", "适合：有点击、有花费，但订单少，想判断是不是 listing 或竞争力问题。", "用法")
    dark_prompt(slide, 0.95, 1.85, 5.35, 2.05, "复制提示词", "我怀疑这个 SKU 不是没流量，而是点击后不转化。请帮我找支持证据、反证、还缺什么证据，以及下一步最小检查动作。")
    items = [("支持证据", "点击、CVR、价格、评价、竞品对照"), ("反证", "样本太少、流量词不准、时间窗不对"), ("下一步", "先补页面/竞品/广告哪类证据")]
    for i, (t, b) in enumerate(items):
        small_card(slide, 7.0, 1.85 + i * 1.0, 5.35, 0.78, t, b, [COLORS["green_soft"], COLORS["red_soft"], COLORS["orange_soft"]][i], [COLORS["green"], COLORS["red"], COLORS["orange"]][i])
    rect(slide, 2.0, 5.85, 9.35, 0.55, fill="111827", line=None, radius=True)
    text_box(slide, "关键：不是让它替你拍脑袋，而是让它先把证据找齐。", 2.3, 6.0, 8.75, 0.16, size=12.5, color="FFFFFF", bold=True, align="center")


def ad_read(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    title(slide, 15, "用法 3：广告页面只读归类", "适合：你已经打开广告页面，想先分清问题属于哪一类。", "用法")
    browser_mock(slide, 0.95, 1.82, 5.4, 3.65, "广告页面字段", ["曝光 / 点击", "花费 / CPC", "订单 / 销售", "CTR / CVR / ACOS"], active_idx=3)
    dark_prompt(slide, 7.0, 1.82, 5.35, 1.78, "复制提示词", "只读看这个广告页面，先判断问题更像曝光不足、点击不足、点击后不转化，还是效率问题。每个判断后面写证据。")
    right = rect(slide, 7.0, 4.0, 5.35, 1.45, fill="FFFFFF", line="D7DEE9", radius=True)
    set_shape_text(right, "不要一上来让它调广告", "先让它归类问题，再生成待确认动作清单。", title_size=15, body_size=12.5, title_color=COLORS["red"], body_color="334155", margin=0.2)


def action_list(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    title(slide, 16, "用法 4：动作清单，不是直接执行", "适合：你想让 Codex 帮你整理可操作建议，但还没准备让它改系统。", "用法")
    dark_prompt(slide, 0.95, 1.85, 5.35, 2.05, "复制提示词", "按规则生成待确认动作清单，不要直接执行。清单里写清楚：对象、建议动作、依据、风险、执行后怎么回读验证。")
    table = rect(slide, 7.0, 1.85, 5.35, 3.55, fill="FFFFFF", line="D7DEE9", radius=True)
    set_shape_text(table, "清单必须包含", "对象\n建议动作\n依据\n风险\n回读验证方式", title_size=17, body_size=15, title_color=COLORS["blue"], body_color="334155", margin=0.28)
    rect(slide, 2.0, 5.92, 9.35, 0.55, fill="111827", line=None, radius=True)
    text_box(slide, "看到清单后，人确认，再谈执行。", 2.3, 6.07, 8.75, 0.16, size=12.8, color="FFFFFF", bold=True, align="center")


def advanced_ops(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    title(slide, 17, "再往前一步：重复运营动作，可以搬到 Codex", "这一页先不细讲，只是让你知道它的上限。今天先会用，后面再把流程做深。", "高阶")
    left = rect(slide, 0.95, 1.8, 5.1, 3.55, fill="FFFFFF", line="D7DEE9", radius=True)
    set_shape_text(
        left,
        "今天先做到",
        "只读当前页面\n整理证据\n生成待确认动作清单\n让人检查后再执行",
        title_size=18,
        body_size=15.5,
        title_color=COLORS["blue"],
        body_color="334155",
        margin=0.28,
    )
    right = rect(slide, 7.25, 1.8, 5.1, 3.55, fill="111827", line=None, radius=True)
    set_shape_text(
        right,
        "用熟以后可以变成",
        "在 Codex 里运营一部分重复工作\n不是每次都切回广告系统、库存系统",
        title_size=18,
        body_size=15.2,
        title_color="FFFFFF",
        body_color="E5E7EB",
        margin=0.28,
    )
    line(slide, 6.18, 3.55, 7.08, 3.55, color=COLORS["orange"], width=2.2)
    cases = [
        ("写文案", "按产品信息生成标题、五点、卖点草稿"),
        ("低效词处理", "按规则识别、生成调整建议，确认后执行"),
        ("高效词跟进", "发现有效词，持续放量或沉淀打法"),
        ("滞销跟进", "按库存、利润、广告、市场证据生成跟进动作"),
    ]
    fills = [COLORS["blue_soft"], COLORS["orange_soft"], COLORS["green_soft"], COLORS["purple_soft"]]
    accents = [COLORS["blue"], COLORS["orange"], COLORS["green"], COLORS["purple"]]
    for i, (t, b) in enumerate(cases):
        small_card(slide, 0.95 + i * 3.08, 5.45, 2.65, 0.95, t, b, fills[i], accents[i])
    text_box(slide, "关键不是“AI 会什么”，而是：凡是你反复做、规则清楚、结果可检查的事，都可以逐步交给它跑。", 1.45, 6.73, 10.45, 0.24, size=12.4, color=COLORS["ink"], bold=True, align="center")


def repair(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    title(slide, 18, "如果它答得不好，直接这样追问", "不要重新来一遍。用追问把它拉回证据和边界。", "追问")
    prompts = [("太空泛", "把每个结论后面的证据字段列出来。"), ("乱下结论", "把支持证据和反证分开写。"), ("想直接改", "先生成待确认动作清单，不要执行。"), ("看不懂字段", "把你不理解的业务词列出来，我逐个确认。"), ("没有下一步", "给我一个最小检查动作，不要给大而全建议。"), ("缺少边界", "重写一次，明确哪些结论只基于当前页面。")]
    for i, (t, b) in enumerate(prompts):
        x = 0.88 + (i % 2) * 6.1
        y = 1.82 + (i // 2) * 1.22
        small_card(slide, x, y, 5.3, 0.88, t, b, COLORS["orange_soft"], COLORS["orange"])


def boundaries(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    title(slide, 19, "四条边界，先记住", "能不能放心用，关键看边界是否清楚。", "边界")
    items = [("只读优先", "新手先不让它修改系统"), ("不越权", "只看你有权限或你提供的内容"), ("不迷信", "结论要能回到证据核对"), ("不装懂", "内部规则不确定就标待确认")]
    for i, (t, b) in enumerate(items):
        x = 0.95 + (i % 2) * 6.08
        y = 2.0 + (i // 2) * 1.55
        small_card(slide, x, y, 5.35, 1.05, t, b, [COLORS["blue_soft"], COLORS["red_soft"], COLORS["green_soft"], COLORS["purple_soft"]][i], [COLORS["blue"], COLORS["red"], COLORS["green"], COLORS["purple"]][i])
    rect(slide, 2.2, 5.75, 8.95, 0.62, fill="111827", line=None, radius=True)
    text_box(slide, "让它先跑第一遍，不是让它替你担责任。", 2.5, 5.93, 8.35, 0.16, size=13, color="FFFFFF", bold=True, align="center")


def summary(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    title(slide, 20, "从不知道 Codex，到第一次用起来", "按这个顺序走就行。", "总览")
    left = rect(slide, 0.95, 1.8, 5.35, 4.25, fill="FFFFFF", line="D7DEE9", radius=True)
    set_shape_text(left, "顺序", "1. 知道 Codex 是工作流助手\n2. 安装公司版 Codex\n3. 放入 amazon-sales-starter 技能包\n4. 打开一个业务页面\n5. 复制只读提示词\n6. 看证据、风险和下一步", title_size=18, body_size=14.3, title_color=COLORS["blue"], body_color="334155", margin=0.28)
    right = rect(slide, 7.0, 1.8, 5.35, 4.25, fill="111827", line=None, radius=True)
    set_shape_text(right, "第一句", "只读帮我看当前页面。\n按 Amazon 销售入门技能包，先识别页面类型、可见字段、初步判断、风险和下一步建议。\n不要直接修改任何内容。", title_size=18, body_size=14, title_color="FFFFFF", body_color="E5E7EB", margin=0.28)


SLIDES = [
    cover,
    what_is_codex,
    chatgpt_difference,
    why_sales_care,
    blank_codex,
    sales_skill_pack,
    after_skill,
    install,
    first_three_steps,
    first_prompt,
    expected_output,
    current_page,
    competitor,
    conversion,
    ad_read,
    action_list,
    advanced_ops,
    repair,
    boundaries,
    summary,
]


def write_quick_card():
    QUICK_CARD.write_text(
        """# 销售同事 Codex 入门速查

## Codex 是什么

Codex 是一个能在你授权范围内，帮你看页面、读文件、整理证据、生成清单的 AI 工作助手。

## 第一次使用顺序

1. 安装公司版 Codex。
2. 放入 `amazon-sales-starter` 销售技能包。
3. 打开一个你平时看的业务页面。
4. 复制下面这句话给 Codex：

```text
只读帮我看当前页面。
按 Amazon 销售入门技能包，先识别页面类型、可见字段、初步判断、风险和下一步建议。
不要直接修改任何内容。
```

## 检查它的输出

- 有没有说清楚页面类型？
- 有没有列出看到的字段？
- 判断后面有没有证据？
- 有没有标出风险或不确定项？
- 有没有保持只读、不直接修改？

## 常用追问

```text
把每个结论后面的证据字段列出来。
```

```text
把支持证据和反证分开写。
```

```text
先生成待确认动作清单，不要执行。
```
""",
        encoding="utf-8",
    )


def build():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    pres = Presentation()
    pres.slide_width = Inches(SLIDE_W)
    pres.slide_height = Inches(SLIDE_H)
    pres.core_properties.title = "销售同事第一次认识 Codex"
    pres.core_properties.subject = "零认知同事版 V5"
    pres.core_properties.author = "Codex"
    for fn in SLIDES:
        fn(pres)
    pres.save(PPTX_PATH)
    write_quick_card()
    print(PPTX_PATH)
    print(QUICK_CARD)


if __name__ == "__main__":
    build()
