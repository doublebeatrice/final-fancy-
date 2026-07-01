from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from generate_codex_sales_share_v2 import (
    COLORS,
    SLIDE_H,
    SLIDE_W,
    background,
    browser_mock,
    card,
    footer,
    header,
    line,
    pill,
    rect,
    set_shape_text,
    small_card,
    step_circle,
    text_box,
)


OUT_DIR = Path(r"D:\ad-ops-workbench\outputs\codex_sales_share_v3")
PPTX_PATH = OUT_DIR / "sales_codex_share_v3_skill_ready.pptx"
NOTES_PATH = OUT_DIR / "sales_codex_share_v3_speaker_script.md"


def cover(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    rect(slide, 0, 0, 5.35, SLIDE_H, fill="111827", line=None, radius=False)
    text_box(slide, "别让每个人", 0.72, 1.18, 4.05, 0.52, size=29, color="FFFFFF", bold=True)
    text_box(slide, "从零教 AI", 0.72, 1.88, 4.05, 0.52, size=32, color="93C5FD", bold=True)
    text_box(slide, "给销售一个装好技能的 Codex", 0.75, 2.78, 3.95, 0.28, size=14, color="E5E7EB")
    text_box(slide, "今天讲的不是“AI 很强”，而是：怎么让刚装好的 Codex 先懂销售最大公约数，开箱就能帮你只读看页面。", 0.75, 3.38, 3.8, 0.9, size=13.5, color="CBD5E1")
    pill(slide, "销售技能包", 0.75, 5.22, 1.45, fill="EAF1FF", color="1D4ED8")
    pill(slide, "开箱只读", 2.38, 5.22, 1.28, fill="E6F7F4", color="0F766E")

    rect(slide, 6.25, 0.9, 5.35, 1.18, fill="FFFFFF", line="D7DEE9", radius=True)
    text_box(slide, "裸 Codex", 6.58, 1.24, 1.6, 0.25, size=15, color=COLORS["red"], bold=True)
    text_box(slide, "不知道你的职责、业务词、页面、判断习惯。", 8.0, 1.25, 3.1, 0.22, size=11.5, color=COLORS["muted"])
    line(slide, 8.85, 2.25, 8.85, 2.75, color=COLORS["orange"], width=2)
    rect(slide, 6.25, 2.9, 5.35, 1.18, fill="FFFFFF", line="D7DEE9", radius=True)
    text_box(slide, "装技能包", 6.58, 3.24, 1.7, 0.25, size=15, color=COLORS["blue"], bold=True)
    text_box(slide, "先懂 Amazon 销售常见术语和基本判断链。", 8.0, 3.25, 3.15, 0.22, size=11.5, color=COLORS["muted"])
    line(slide, 8.85, 4.25, 8.85, 4.75, color=COLORS["orange"], width=2)
    rect(slide, 6.25, 4.9, 5.35, 1.18, fill="111827", line=None, radius=True)
    text_box(slide, "开箱第一句", 6.58, 5.24, 1.8, 0.25, size=15, color="FFFFFF", bold=True)
    text_box(slide, "只读帮我看当前页面。", 8.0, 5.25, 2.8, 0.22, size=12.5, color="E5E7EB", bold=True)
    footer(slide, 1)


def blank_truth(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    header(slide, 2, "先说清楚：刚装好的 Codex 是空白的", "它不会天然懂 Amazon 销售，也不会天然知道公司内部字段。不能把开箱即用讲成“它天生都懂”。", "真相")
    left = rect(slide, 0.95, 2.0, 5.35, 3.6, fill="FFFFFF", line="D7DEE9", radius=True)
    set_shape_text(left, "裸 Codex 不知道", "你负责什么品类\nSKU / ASIN 是什么关系\n滞销意味着什么\n广告字段怎么判断\n哪些动作不能直接做", title_size=17, body_size=15, title_color=COLORS["red"], body_color="334155", margin=0.3)
    right = rect(slide, 7.02, 2.0, 5.35, 3.6, fill="FFFFFF", line="D7DEE9", radius=True)
    set_shape_text(right, "所以不能这样开场", "“你直接问它 SKU 就行”\n“它会自己懂业务”\n“你慢慢教它就好”\n\n这会让同事觉得麻烦。", title_size=17, body_size=15, title_color=COLORS["orange"], body_color="334155", margin=0.3)
    text_box(slide, "正确目标：不要让每个销售都从零带新人。", 2.2, 6.05, 8.9, 0.34, size=15, color=COLORS["ink"], bold=True, align="center")


def skill_solution(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    header(slide, 3, "解决方案：先给它装一个销售通用技能包", "销售职责有大量最大公约数。我们先把这些共识装进去，同事不用每个人重新解释一遍。", "方案")
    rect(slide, 0.95, 1.95, 3.3, 3.75, fill="FFFFFF", line="D7DEE9", radius=True)
    text_box(slide, "过去", 1.25, 2.35, 1.0, 0.3, size=16, color=COLORS["red"], bold=True)
    text_box(slide, "每个人第一次用\n都要解释一堆：\n职责、术语、页面、判断逻辑。", 1.25, 3.0, 2.3, 1.1, size=15.5, color="334155")
    line(slide, 4.55, 3.75, 5.2, 3.75, color=COLORS["orange"], width=2.2)
    rect(slide, 5.45, 1.65, 2.45, 4.35, fill="111827", line=None, radius=True)
    text_box(slide, "销售\n技能包", 5.85, 2.65, 1.65, 0.72, size=24, color="FFFFFF", bold=True, align="center")
    text_box(slide, "职责 / 术语 / 页面 / 判断链 / 边界", 5.8, 4.18, 1.75, 0.55, size=11.5, color="CBD5E1", align="center")
    line(slide, 8.18, 3.75, 8.83, 3.75, color=COLORS["orange"], width=2.2)
    rect(slide, 9.1, 1.95, 3.3, 3.75, fill="FFFFFF", line="D7DEE9", radius=True)
    text_box(slide, "现在", 9.4, 2.35, 1.0, 0.3, size=16, color=COLORS["green"], bold=True)
    text_box(slide, "同事只要打开页面\n说一句只读任务\nCodex 先按销售共识跑第一遍。", 9.4, 3.0, 2.35, 1.1, size=15.5, color="334155")
    text_box(slide, "这才是“开箱即用”：不是它天生懂，而是先给它装好销售手册。", 1.7, 6.25, 10.0, 0.3, size=14, color=COLORS["ink"], bold=True, align="center")


def skill_contents(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    header(slide, 4, "销售技能包里先装什么？只装最大公约数", "不要一开始追求全公司所有细节。先让它懂销售同事每天都会遇到的基础上下文。", "内容")
    items = [
        ("销售职责", "产品、广告、库存、利润、竞品、前台表现", COLORS["blue_soft"], COLORS["blue"]),
        ("常见术语", "SKU、ASIN、CTR、CVR、ACOS、FBA、滞销", COLORS["teal_soft"], COLORS["teal"]),
        ("常见页面", "广告系统、库存产品系统、SIF、Amazon 前台", COLORS["orange_soft"], COLORS["orange"]),
        ("判断方法", "先看目标，再看产品能否接流量，再看广告问题", COLORS["purple_soft"], COLORS["purple"]),
        ("安全边界", "默认只读、不越权、不乱改、动作先列清单", COLORS["green_soft"], COLORS["green"]),
    ]
    for i, (title, body, fill, accent) in enumerate(items):
        x = 0.95 + (i % 2) * 6.05
        y = 1.85 + (i // 2) * 1.28
        w = 5.25 if i < 4 else 11.3
        if i == 4:
            x = 0.95
        small_card(slide, x, y, w, 0.88, title, body, fill, accent)
    rect(slide, 2.1, 6.1, 9.15, 0.52, fill="111827", line=None, radius=True)
    text_box(slide, "技能包不是替代业务判断，而是让 Codex 不再从零开始。", 2.35, 6.24, 8.65, 0.16, size=12.5, color="FFFFFF", bold=True, align="center")


def first_prompt(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    header(slide, 5, "装好技能包后，第一句话就可以很简单", "同事不用先讲一堆背景。先让 Codex 只读当前页面，按销售技能包识别信息和缺口。", "开箱")
    browser_mock(slide, 0.95, 1.88, 5.4, 3.65, "当前 Chrome 页面", ["页面类型", "可见对象：SKU / ASIN / 搜索词", "关键字段", "不理解的字段"], active_idx=2)
    prompt = rect(slide, 7.0, 1.88, 5.35, 3.65, fill="111827", line=None, radius=True)
    set_shape_text(prompt, "直接复制这句话", "只读帮我看当前页面。\n按 Amazon 销售入门技能包，先识别页面类型、可见字段、初步判断、风险和下一步建议。\n不要直接修改任何内容。", title_size=17, body_size=15, title_color="FFFFFF", body_color="E5E7EB", margin=0.28)
    text_box(slide, "关键变化：同事不是在带新人，而是在叫一个已经读过销售手册的助手先跑一遍。", 1.2, 6.14, 10.9, 0.3, size=14, color=COLORS["ink"], bold=True, align="center")


def effect_demo(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    header(slide, 6, "现场效果要演示“技能包加持后，它先懂什么”", "不是演示一个复杂闭环，而是让大家看到：它先会识别页面、解释字段、列证据、提醒边界。", "演示")
    steps = [
        ("1", "打开页面", "销售平时看的 Chrome 页面"),
        ("2", "说一句话", "只读帮我看当前页面"),
        ("3", "它先识别", "页面类型、字段、对象、时间窗"),
        ("4", "它给草稿", "看到的信息、判断、风险、下一步"),
    ]
    for i, (num, title, body) in enumerate(steps):
        x = 0.95 + i * 3.05
        box = rect(slide, x, 2.18, 2.45, 2.72, fill="FFFFFF", line="D7DEE9", radius=True)
        step_circle(slide, x + 0.25, 2.48, num, [COLORS["blue"], COLORS["orange"], COLORS["green"], COLORS["purple"]][i])
        text_box(slide, title, x + 0.25, 3.22, 1.9, 0.26, size=15, color=COLORS["ink"], bold=True)
        text_box(slide, body, x + 0.25, 3.72, 1.88, 0.46, size=12, color=COLORS["muted"])
        if i < 3:
            line(slide, x + 2.55, 3.52, x + 2.92, 3.52, color="CBD5E1", width=1.3)
    rect(slide, 2.0, 5.75, 9.35, 0.6, fill="FFFFFF", line="D7DEE9", radius=True)
    text_box(slide, "演示成功标准：大家觉得“我打开自己的页面也能试一句”。", 2.25, 5.93, 8.85, 0.16, size=13, color=COLORS["blue"], bold=True, align="center")


def scenario_competitor(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    header(slide, 7, "场景一：竞品市场，不让它只看一个链接", "销售技能包会提醒它从价格带、评论门槛、主图打法、卖点结构去看一片市场。", "场景")
    blocks = [
        ("输入", "搜索词\n或 10 个竞品链接", COLORS["blue"]),
        ("技能包提醒", "不要只总结单品\n要看市场结构", COLORS["orange"]),
        ("输出", "价格带 / 评论门槛 / 主图打法 / 卖点差异 / 我们差距", COLORS["green"]),
    ]
    for i, (title, body, accent) in enumerate(blocks):
        x = 0.95 + i * 4.1
        box = rect(slide, x, 2.0, 3.35, 3.25, fill="FFFFFF", line="D7DEE9", radius=True)
        set_shape_text(box, title, body, title_size=17, body_size=15, title_color=accent, body_color="334155", margin=0.28)
        if i < 2:
            line(slide, x + 3.55, 3.58, x + 3.88, 3.58, color="CBD5E1", width=1.5)
    prompt = rect(slide, 1.25, 5.82, 10.82, 0.56, fill="FFFFFF", line="D7DEE9", radius=True)
    set_shape_text(prompt, "提示词：帮我看这个搜索词下的 10 个竞品，整理价格、评论、主图和卖点差异。", title_size=12.5, title_color=COLORS["ink"], title_bold=True, align="center", valign="middle", margin=0.04)


def scenario_conversion(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    header(slide, 8, "场景二：转化问题，让它找证据和反证", "技能包不会让它直接说“优化 listing”。它会先判断问题在曝光、点击、转化还是效率。", "场景")
    quote = rect(slide, 0.95, 2.05, 3.35, 3.45, fill="111827", line=None, radius=True)
    set_shape_text(quote, "人的判断", "我怀疑这个产品不是没流量，而是点击后不转化。", title_size=15, body_size=18, title_color="FFFFFF", body_color="E5E7EB", margin=0.28)
    items = [
        ("先看曝光", "有没有流量覆盖", COLORS["blue_soft"], COLORS["blue"]),
        ("再看点击", "CTR、图片、标题、价格", COLORS["teal_soft"], COLORS["teal"]),
        ("再看转化", "CVR、评价、价格、竞品", COLORS["orange_soft"], COLORS["orange"]),
        ("最后给反证", "也可能是词不准或样本太少", COLORS["red_soft"], COLORS["red"]),
    ]
    for i, (title, body, fill, accent) in enumerate(items):
        x = 5.0 + (i % 2) * 3.55
        y = 2.05 + (i // 2) * 1.35
        small_card(slide, x, y, 3.05, 0.94, title, body, fill, accent)
    rect(slide, 5.0, 5.15, 6.62, 0.74, fill="FFFFFF", line="D7DEE9", radius=True)
    text_box(slide, "输出：支持证据 / 反证 / 还缺什么 / 下一步最小检查动作", 5.28, 5.38, 6.05, 0.16, size=12.5, color=COLORS["blue"], bold=True, align="center")


def scenario_ad_readonly(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    header(slide, 9, "场景三：广告只读诊断，先分清是哪类问题", "同事不需要先教它所有广告字段。技能包里已经有 CTR、CVR、ACOS、曝光、点击、花费的基础解释。", "场景")
    browser_mock(slide, 0.95, 1.92, 5.55, 3.65, "广告页面可见字段", ["曝光 / 点击 / 花费", "CTR / CPC", "订单 / 销售", "ACOS / CVR"], active_idx=3)
    right = rect(slide, 7.0, 1.92, 5.35, 3.65, fill="FFFFFF", line="D7DEE9", radius=True)
    set_shape_text(right, "它先归类", "曝光不足\n点击不足\n点击后不转化\n有单但效率差\n样本太少先不下结论", title_size=17, body_size=15, title_color=COLORS["blue"], body_color="334155", margin=0.28)
    text_box(slide, "重点：先让它只读归类，不要一上来就让它调广告。", 1.6, 6.12, 10.1, 0.28, size=14, color=COLORS["ink"], bold=True, align="center")


def action_gate(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    header(slide, 10, "最后一公里可以做，但技能包默认先拦住", "这个边界要讲清楚：不是不让自动化，而是先清单、先确认、再执行、再回读。", "动作")
    items = [
        ("先只读", "页面和证据先看清楚", COLORS["blue_soft"], COLORS["blue"]),
        ("列清单", "对象、动作、依据、风险", COLORS["orange_soft"], COLORS["orange"]),
        ("人确认", "范围、强度、例外项", COLORS["purple_soft"], COLORS["purple"]),
        ("回读验证", "确认动作是否真的落地", COLORS["green_soft"], COLORS["green"]),
    ]
    for i, (title, body, fill, accent) in enumerate(items):
        x = 0.92 + i * 3.07
        box = rect(slide, x, 2.15, 2.45, 2.75, fill="FFFFFF", line="D7DEE9", radius=True)
        set_shape_text(box, title, body, title_size=16, body_size=12.5, title_color=accent, body_color="475569", margin=0.22)
        if i < 3:
            line(slide, x + 2.55, 3.52, x + 2.94, 3.52, color="CBD5E1", width=1.4)
    banner = rect(slide, 1.75, 5.75, 9.85, 0.62, fill="111827", line=None, radius=True)
    set_shape_text(banner, "技能包的默认规则：没有授权，不直接改。", title_size=14, title_color="FFFFFF", title_bold=True, align="center", valign="middle", margin=0.04)


def prompts(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    header(slide, 11, "发给同事的不是一堆教程，而是这几句开箱提示词", "技能包装好后，提示词可以非常短，因为基础背景已经在技能里了。", "提示词")
    data = [
        ("当前页面", "只读帮我看当前页面，按销售技能包整理看到的信息、风险和下一步建议。"),
        ("广告诊断", "只读看这个广告页面，先判断问题更像曝光、点击、转化还是效率。"),
        ("竞品市场", "帮我看这个搜索词下的 10 个竞品，整理价格、评论、主图和卖点差异。"),
        ("转化假设", "我怀疑点击后不转化，帮我找支持证据、反证和还缺什么。"),
        ("动作清单", "按规则生成待确认动作清单，不要直接执行。"),
        ("复盘草稿", "把这次调整前后的证据整理成复盘草稿，标出还缺什么。"),
    ]
    fills = [COLORS["blue_soft"], COLORS["teal_soft"], COLORS["orange_soft"], COLORS["purple_soft"], COLORS["green_soft"], COLORS["red_soft"]]
    accents = [COLORS["blue"], COLORS["teal"], COLORS["orange"], COLORS["purple"], COLORS["green"], COLORS["red"]]
    for i, (title, body) in enumerate(data):
        x = 0.88 + (i % 2) * 6.1
        y = 1.85 + (i // 2) * 1.45
        small_card(slide, x, y, 5.3, 1.02, title, body, fills[i], accents[i])


def install_skill(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    header(slide, 12, "会后路径：安装 Codex，再装销售技能包", "安装本身不要在开场讲。先看效果，想试的人再按公司指南和技能包文件走。", "会后")
    steps = [
        ("1", "安装公司版 Codex", "按企业微信《Codex安装&使用教程》。"),
        ("2", "拿到销售技能包", "文件夹：amazon-sales-starter。"),
        ("3", "放进 skills 目录", "或按 IT/助理给的方式导入。"),
        ("4", "打开 Chrome 页面", "先从只读当前页开始。"),
    ]
    for i, (num, title, body) in enumerate(steps):
        x = 0.95 + (i % 2) * 6.08
        y = 1.95 + (i // 2) * 1.35
        box = rect(slide, x, y, 5.35, 0.95, fill="FFFFFF", line="D7DEE9", radius=True)
        step_circle(slide, x + 0.22, y + 0.25, num, [COLORS["blue"], COLORS["orange"], COLORS["green"], COLORS["purple"]][i])
        text_box(slide, title, x + 0.88, y + 0.18, 3.8, 0.22, size=14.2, color=COLORS["ink"], bold=True)
        text_box(slide, body, x + 0.88, y + 0.52, 4.1, 0.2, size=10.8, color=COLORS["muted"])
    rect(slide, 2.0, 5.35, 9.35, 0.72, fill="111827", line=None, radius=True)
    text_box(slide, "卡住不用硬扛：让助理远程协助安装 Codex 和技能包。", 2.25, 5.58, 8.85, 0.16, size=13.5, color="FFFFFF", bold=True, align="center")


def practice(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    header(slide, 13, "现场练习：不是教它业务，而是测试技能包是否够用", "每个人选一个当前页面，跑一句只读提示词，看它能不能先识别、先归类、先列问题。", "练习")
    blocks = [
        ("打开页面", "广告 / 库存 / Amazon 前台 / 竞品搜索结果", COLORS["blue_soft"], COLORS["blue"]),
        ("说一句话", "只读帮我看当前页面，按销售技能包整理。", COLORS["orange_soft"], COLORS["orange"]),
        ("看三件事", "识别准不准\n字段懂不懂\n建议能不能查", COLORS["green_soft"], COLORS["green"]),
    ]
    for i, (title, body, fill, accent) in enumerate(blocks):
        x = 0.95 + i * 4.1
        box = rect(slide, x, 2.05, 3.35, 3.35, fill="FFFFFF", line="D7DEE9", radius=True)
        set_shape_text(box, title, body, title_size=17, body_size=15, title_color=accent, body_color="334155", margin=0.28)
        if i < 2:
            line(slide, x + 3.55, 3.68, x + 3.88, 3.68, color="CBD5E1", width=1.5)
    text_box(slide, "如果它问的问题更少、输出更像销售语言，技能包就发挥作用了。", 1.65, 6.05, 10.0, 0.3, size=14, color=COLORS["ink"], bold=True, align="center")


def safety(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    header(slide, 14, "放心试，但技能包也要写清四条边界", "越是想推广，越要把边界写进默认规则里。", "边界")
    items = [
        ("不越权", "只看本人能访问或用户提供的内容", COLORS["blue_soft"], COLORS["blue"]),
        ("不乱改", "默认只读；修改必须有确认清单", COLORS["red_soft"], COLORS["red"]),
        ("不迷信", "结论后面要有证据，能回到页面核对", COLORS["orange_soft"], COLORS["orange"]),
        ("不装懂", "内部特殊规则缺失时，先标待确认", COLORS["teal_soft"], COLORS["teal"]),
    ]
    for i, (title, body, fill, accent) in enumerate(items):
        x = 0.95 + (i % 2) * 6.08
        y = 2.0 + (i // 2) * 1.55
        box = rect(slide, x, y, 5.35, 1.15, fill=fill, line=None, radius=True)
        set_shape_text(box, title, body, title_size=15, body_size=11.5, title_color=accent, body_color="334155", margin=0.18)
    rect(slide, 2.2, 5.7, 8.95, 0.68, fill="FFFFFF", line="D7DEE9", radius=True)
    text_box(slide, "一句话：技能包让它先懂基础，不代表让它替你担责任。", 2.5, 5.9, 8.35, 0.2, size=13.2, color=COLORS["ink"], bold=True, align="center")


def closing(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    background(slide)
    rect(slide, 0.8, 0.78, 11.75, 5.55, fill="111827", line=None, radius=True)
    text_box(slide, "开箱即用，不是空白上手", 1.35, 1.48, 10.7, 0.55, size=29, color="FFFFFF", bold=True, align="center")
    text_box(slide, "是先给 Codex 装好销售技能，再让它只读跑第一遍。", 1.75, 2.3, 9.9, 0.36, size=18, color="E5E7EB", align="center")
    labels = [("装技能", COLORS["blue"]), ("打开页面", COLORS["orange"]), ("只读一句", COLORS["green"])]
    for i, (txt, accent) in enumerate(labels):
        box = rect(slide, 2.05 + i * 3.2, 3.55, 2.35, 1.0, fill="FFFFFF", line=None, radius=True)
        set_shape_text(box, txt, None, title_size=17, title_color=accent, title_bold=True, align="center", valign="middle", margin=0.04)
    text_box(slide, "今天的行动：拿到技能包，找一个当前页面，试一次只读。", 1.55, 5.25, 10.25, 0.28, size=13.5, color="CBD5E1", align="center")
    footer(slide, 15)


SLIDES = [
    cover,
    blank_truth,
    skill_solution,
    skill_contents,
    first_prompt,
    effect_demo,
    scenario_competitor,
    scenario_conversion,
    scenario_ad_readonly,
    action_gate,
    prompts,
    install_skill,
    practice,
    safety,
    closing,
]

NOTES = [
    ("封面", "开场强调：今天不是让大家从零教 AI，而是给销售一个装好技能的 Codex。"),
    ("裸 Codex 是空白", "先承认真实问题，避免过度承诺开箱就懂全部业务。"),
    ("销售技能包", "提出解决方案：把销售最大公约数沉淀成技能包。"),
    ("技能包内容", "解释技能包只装通用共识，不装个人经验和敏感路径。"),
    ("第一句话", "给出最重要的开箱提示词：只读当前页面。"),
    ("现场效果", "演示目标是证明它能识别页面和字段，不是证明全自动闭环。"),
    ("竞品市场", "场景要从单链接升级成市场结构分析。"),
    ("转化问题", "用证据和反证体现它不是拍脑袋。"),
    ("广告诊断", "展示技能包内置广告指标解释，先做问题归类。"),
    ("动作闸门", "说明最后一公里可以做，但先清单、先确认、再回读。"),
    ("提示词", "让同事拍照或会后复制，降低入门门槛。"),
    ("会后路径", "安装 Codex 之后再安装销售技能包。"),
    ("练习", "现场练习是测试技能包，而不是让同事从零教学。"),
    ("边界", "让组长也能接受：不越权、不乱改、不迷信、不装懂。"),
    ("收尾", "行动口径：装技能包，打开页面，只读试一次。"),
]


def build():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    pres = Presentation()
    pres.slide_width = Inches(SLIDE_W)
    pres.slide_height = Inches(SLIDE_H)
    pres.core_properties.title = "销售版 Codex：技能包开箱即用"
    pres.core_properties.subject = "正式分享课件 V3"
    pres.core_properties.author = "Codex"
    for fn in SLIDES:
        fn(pres)
    pres.save(PPTX_PATH)

    lines = ["# 销售版 Codex 分享 V3 讲稿要点", ""]
    for idx, (title, body) in enumerate(NOTES, start=1):
        lines.append(f"## Slide {idx:02d} - {title}")
        lines.append(body)
        lines.append("")
    NOTES_PATH.write_text("\n".join(lines), encoding="utf-8")
    print(PPTX_PATH)
    print(NOTES_PATH)


if __name__ == "__main__":
    build()
