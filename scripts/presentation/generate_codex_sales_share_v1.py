from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


OUT_DIR = Path(r"D:\ad-ops-workbench\outputs\codex_sales_share_v1")
PPTX_PATH = OUT_DIR / "codex_sales_share_v1.pptx"
NOTES_PATH = OUT_DIR / "codex_sales_share_v1_speaker_notes.md"

W, H = 13.333, 7.5

COLORS = {
    "bg": "F7F8FA",
    "paper": "FFFFFF",
    "ink": "111827",
    "muted": "64748B",
    "line": "D6DDE8",
    "blue": "2563EB",
    "blue_dark": "1E3A8A",
    "teal": "0F766E",
    "green": "10B981",
    "orange": "F97316",
    "red": "EF4444",
    "purple": "7C3AED",
    "yellow": "FBBF24",
    "soft_blue": "EAF1FF",
    "soft_teal": "E7F7F4",
    "soft_orange": "FFF3E8",
    "soft_green": "EAFBF4",
    "soft_red": "FEECEC",
    "soft_purple": "F3EAFE",
}

FONT = "Microsoft YaHei"
FONT_LATIN = "Arial"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color.replace("#", ""))


def set_run_font(run, size=18, color="111827", bold=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    rpr = run._r.get_or_add_rPr()
    rfonts = rpr.find(qn("a:rFonts"))
    if rfonts is None:
        rfonts = OxmlElement("a:rFonts")
        rpr.append(rfonts)
    rfonts.set(qn("a:latin"), FONT_LATIN)
    rfonts.set(qn("a:ea"), FONT)


def add_text(slide, text, x, y, w, h, size=20, color="111827", bold=False,
             align="left", valign="top", margin=0.03, line_spacing=1.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    if valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    p = tf.paragraphs[0]
    p.alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }[align]
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    set_run_font(run, size=size, color=color, bold=bold)
    return box


def add_multiline(slide, lines, x, y, w, h, size=18, color="111827",
                  bullet=False, gap=8, bold_first=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        if bullet:
            p.text = ""
            p._p.get_or_add_pPr().insert(0, p._p.get_or_add_pPr())
        run = p.add_run()
        run.text = line
        set_run_font(run, size=size, color=color, bold=(bold_first and idx == 0))
    return box


def add_rect(slide, x, y, w, h, fill="FFFFFF", line=None, radius=False,
             transparency=0, shadow=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.fill.transparency = transparency
    if line:
        shp.line.color.rgb = rgb(line)
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    if shadow:
        effect = shp.shadow
        effect.inherit = False
        effect.blur_radius = Pt(4)
        effect.distance = Pt(1.2)
        effect.angle = 45
        effect.transparency = 0.82
    return shp


def add_line(slide, x1, y1, x2, y2, color="D6DDE8", width=1.2, begin=False, end=False):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    conn.line.color.rgb = rgb(color)
    conn.line.width = Pt(width)
    if end:
        conn.line.end_arrowhead = True
    if begin:
        conn.line.begin_arrowhead = True
    return conn


def add_chip(slide, text, x, y, w, h=0.34, fill="EAF1FF", color="1E3A8A"):
    add_rect(slide, x, y, w, h, fill=fill, line=None, radius=True)
    return add_text(slide, text, x + 0.08, y + 0.035, w - 0.16, h - 0.04,
                    size=11.5, color=color, bold=True, align="center", valign="middle")


def add_title(slide, title, subtitle=None, section=None):
    if section:
        add_chip(slide, section, 0.55, 0.38, 1.55, fill="EAF1FF", color="1E3A8A")
    add_text(slide, title, 0.55, 0.76, 9.8, 0.55, size=26, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.58, 1.26, 10.6, 0.36, size=12.5, color="64748B")


def add_footer(slide, page, label="销售 Codex 分享讨论稿 v1"):
    add_line(slide, 0.55, 7.05, 12.78, 7.05, color="E3E8F0", width=0.7)
    add_text(slide, label, 0.55, 7.14, 4.5, 0.18, size=8.5, color="94A3B8")
    add_text(slide, f"{page:02d}", 12.25, 7.14, 0.5, 0.18, size=8.5, color="94A3B8", align="right")


def add_background(slide):
    add_rect(slide, 0, 0, W, H, fill=COLORS["bg"])


def add_browser(slide, x, y, w, h, title="已登录 Chrome 页面", rows=None, highlight=None):
    add_rect(slide, x, y, w, h, fill="FFFFFF", line="D6DDE8", radius=True, shadow=True)
    add_rect(slide, x, y, w, 0.45, fill="F1F5F9", line=None, radius=True)
    for i, c in enumerate(["EF4444", "FBBF24", "10B981"]):
        slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.18 + i * 0.22), Inches(y + 0.16), Inches(0.09), Inches(0.09)).fill.solid()
        slide.shapes[-1].fill.fore_color.rgb = rgb(c)
        slide.shapes[-1].line.fill.background()
    add_text(slide, title, x + 0.75, y + 0.12, w - 1.0, 0.22, size=9.5, color="475569", valign="middle")
    rows = rows or ["SKU / ASIN / 搜索词", "广告表现", "库存与利润", "竞品页面"]
    row_y = y + 0.72
    for idx, row in enumerate(rows):
        fill = "EAF1FF" if highlight == idx else "FFFFFF"
        add_rect(slide, x + 0.28, row_y + idx * 0.55, w - 0.56, 0.36, fill=fill, line="E2E8F0", radius=True)
        add_text(slide, row, x + 0.42, row_y + idx * 0.06 + idx * 0.55 + 0.04,
                 w - 0.84, 0.22, size=9.5, color="334155")


def add_panel(slide, title, body, x, y, w, h, accent="2563EB", fill="FFFFFF"):
    add_rect(slide, x, y, w, h, fill=fill, line="D6DDE8", radius=True, shadow=False)
    add_rect(slide, x, y, 0.07, h, fill=accent, line=None, radius=False)
    add_text(slide, title, x + 0.22, y + 0.18, w - 0.4, 0.34, size=15.5, bold=True, color="111827")
    add_text(slide, body, x + 0.22, y + 0.68, w - 0.42, h - 0.85, size=12.4, color="475569")


def add_kpi(slide, num, label, x, y, w, color="2563EB"):
    add_text(slide, num, x, y, w, 0.48, size=27, color=color, bold=True, align="center")
    add_text(slide, label, x, y + 0.5, w, 0.32, size=11.5, color="475569", align="center")


def slide_cover(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    add_background(slide)
    add_rect(slide, 0, 0, W, H, fill="F7F8FA")
    add_rect(slide, 0.5, 0.48, 5.0, 5.9, fill="FFFFFF", line="D6DDE8", radius=True, shadow=True)
    add_browser(slide, 6.15, 0.82, 5.95, 3.0, rows=["HUA0165 / 已登录广告页面", "30天、7天、3天表现", "初步判断与待确认点"], highlight=1)
    add_rect(slide, 7.0, 4.18, 4.2, 1.15, fill="111827", line=None, radius=True)
    add_text(slide, "你说一句", 7.25, 4.39, 1.2, 0.32, size=14, color="FFFFFF", bold=True)
    add_line(slide, 8.45, 4.55, 9.1, 4.55, color="94A3B8", width=1.5)
    add_text(slide, "它去看", 9.16, 4.39, 1.15, 0.32, size=14, color="FFFFFF", bold=True)
    add_line(slide, 10.2, 4.55, 10.78, 4.55, color="94A3B8", width=1.5)
    add_text(slide, "出清单", 10.72, 4.39, 1.1, 0.32, size=14, color="FFFFFF", bold=True)
    add_text(slide, "销售怎么开始用 Codex", 0.85, 1.02, 4.25, 0.92, size=34, bold=True, color="111827")
    add_text(slide, "从一个重复小动作开始", 0.88, 1.96, 4.05, 0.45, size=18, color="2563EB", bold=True)
    add_text(slide, "不是换系统，不是多学一套流程。\n是遇到重复步骤时，先叫它帮你跑一遍。", 0.9, 2.72, 3.9, 1.08, size=16, color="334155")
    add_chip(slide, "讨论稿 v1", 0.92, 4.35, 1.2, fill="FFF3E8", color="C2410C")
    add_chip(slide, "销售培训分享", 2.25, 4.35, 1.6, fill="E7F7F4", color="0F766E")
    add_text(slide, "建议开场：先看效果，再讲怎么做到。", 0.92, 5.08, 4.0, 0.34, size=11.5, color="64748B")
    add_footer(slide, 1)
    return slide


def slide_hook(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    add_background(slide)
    add_title(slide, "开场先放效果：让大家先觉得“我想试”", "不要先讲概念、安装、权限、模型。先让听众看到：它确实能少做一段重复步骤。", "开场")
    add_rect(slide, 0.75, 1.95, 5.35, 3.75, fill="FFFFFF", line="D6DDE8", radius=True, shadow=True)
    add_text(slide, "过去", 1.05, 2.2, 0.85, 0.32, size=16, bold=True, color="EF4444")
    add_text(slide, "手动打开多个页面\n自己找字段\n自己判断和整理\n再写成一段话", 1.05, 2.82, 4.7, 1.6, size=18, color="334155")
    add_text(slide, "问题不是不会做，而是每次都要重做一遍。", 1.05, 5.02, 4.7, 0.34, size=12.5, color="64748B")
    add_rect(slide, 7.15, 1.95, 5.35, 3.75, fill="FFFFFF", line="D6DDE8", radius=True, shadow=True)
    add_text(slide, "现场效果", 7.45, 2.2, 1.4, 0.32, size=16, bold=True, color="10B981")
    add_text(slide, "给一句任务\nCodex 看已登录页面\n输出诊断初稿\n人确认、改规则、再决定", 7.45, 2.82, 4.55, 1.6, size=18, color="334155")
    add_text(slide, "重点不是承诺固定几分钟，而是证明它能跟着工作流走。", 7.45, 5.02, 4.5, 0.38, size=12.5, color="64748B")
    add_line(slide, 6.3, 3.78, 6.95, 3.78, color="F97316", width=2.4)
    add_text(slide, "少掉一段重复动作", 5.08, 3.25, 3.0, 0.38, size=14, color="F97316", bold=True, align="center")
    add_footer(slide, 2)
    return slide


def slide_not_chatgpt(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    add_background(slide)
    add_title(slide, "这不是“又一个聊天框”", "如果只是润色、总结、写汇报，ChatGPT 也能做。Codex 的差异要讲在“它能跟着你的工作环境走”。", "核心差异")
    add_panel(slide, "ChatGPT 更像文案搭子", "你把材料复制进去，它帮你改写、总结、提炼。\n\n适合：汇报、复盘、沟通稿、会议纪要。", 0.85, 2.05, 5.1, 3.65, accent="64748B", fill="FFFFFF")
    add_panel(slide, "Codex 更像工作流助手", "它可以在你授权的环境里，看页面、读文件、按步骤操作。\n\n适合：重复查看、重复判断、重复整理、待确认动作清单。", 7.25, 2.05, 5.1, 3.65, accent="2563EB", fill="FFFFFF")
    add_text(slide, "以前：人给 AI 喂材料", 1.25, 5.95, 4.0, 0.32, size=13, color="64748B", align="center")
    add_text(slide, "现在：AI 先去看你允许它看的东西", 7.55, 5.95, 4.5, 0.32, size=13, color="2563EB", bold=True, align="center")
    add_footer(slide, 3)
    return slide


def slide_chrome_extension(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    add_background(slide)
    add_title(slide, "真正的“王炸”：复用本人 Chrome 登录态", "不讲抓包、不讲后台路径。对销售来说，只需要理解：它看的是你本来就能看的页面。", "能力边界")
    add_browser(slide, 0.85, 1.85, 5.5, 3.35, title="你的 Chrome：广告系统 / SIF / SellerInventory", rows=["已登录状态", "本人账号权限", "页面可见内容", "必要时只读分析"], highlight=2)
    add_rect(slide, 7.0, 1.85, 5.45, 3.35, fill="FFFFFF", line="D6DDE8", radius=True, shadow=True)
    add_text(slide, "讲法要稳", 7.35, 2.15, 1.8, 0.35, size=16, bold=True, color="111827")
    points = [
        "不绕权限：只能看本人已有权限",
        "不额外导出：先读页面可见内容",
        "不乱修改：默认只读，动作必须确认",
        "可复用：广告系统跑通，SIF 登录后同理验证",
    ]
    y = 2.75
    colors = ["2563EB", "0F766E", "F97316", "7C3AED"]
    for i, point in enumerate(points):
        add_rect(slide, 7.35, y + i * 0.55, 0.18, 0.18, fill=colors[i], line=None, radius=True)
        add_text(slide, point, 7.68, y + i * 0.50 - 0.02, 4.3, 0.32, size=13.5, color="334155")
    add_text(slide, "这页是给组长也能接受的安全表达：本人权限、可见页面、先只读、后确认。", 0.9, 5.88, 11.7, 0.36, size=12.5, color="64748B", align="center")
    add_footer(slide, 4)
    return slide


def slide_repetition_filter(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    add_background(slide)
    add_title(slide, "什么时候该想到 Codex？一句话：重复做的事", "不要让大家背场景清单。只给一个判断标准：只要这件事有固定证据和固定方法，就值得先交给它跑一遍。", "判断标准")
    cards = [
        ("每周都做", "周会前看一遍、每个 SKU 都看一遍、每次促销都看一遍", "EAF1FF", "2563EB"),
        ("步骤相似", "打开页面、查字段、比较趋势、整理判断，顺序大致固定", "E7F7F4", "0F766E"),
        ("有方法论", "你心里知道先看什么、后看什么、什么情况下怎么判断", "FFF3E8", "F97316"),
        ("结果能查", "它给出的证据、结论、建议都能回到页面或文件里核对", "F3EAFE", "7C3AED"),
    ]
    x0, y0 = 0.75, 2.05
    for i, (title, body, fill, accent) in enumerate(cards):
        x = x0 + (i % 2) * 6.25
        y = y0 + (i // 2) * 1.65
        add_rect(slide, x, y, 5.55, 1.22, fill="FFFFFF", line="D6DDE8", radius=True, shadow=True)
        add_rect(slide, x + 0.22, y + 0.25, 0.58, 0.58, fill=fill, line=None, radius=True)
        add_text(slide, str(i + 1), x + 0.22, y + 0.35, 0.58, 0.24, size=14, color=accent, bold=True, align="center", valign="middle")
        add_text(slide, title, x + 0.98, y + 0.22, 2.2, 0.3, size=15.5, bold=True, color="111827")
        add_text(slide, body, x + 0.98, y + 0.62, 4.25, 0.38, size=11.7, color="64748B")
    add_rect(slide, 2.05, 5.75, 9.2, 0.58, fill="111827", line=None, radius=True)
    add_text(slide, "讲给同事听：不是“你要学 AI”，而是“你把自己的重复方法说给它听”。", 2.25, 5.88, 8.8, 0.25, size=13, color="FFFFFF", bold=True, align="center")
    add_footer(slide, 5)
    return slide


def slide_competitor(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    add_background(slide)
    add_title(slide, "场景 1：竞品和市场，不看一个链接，要看一片", "这个方向容易让销售觉得值：不是帮我看一个竞品，而是帮我把一个搜索词下的 10 个竞品先整理出来。", "场景")
    add_rect(slide, 0.82, 1.9, 3.05, 3.95, fill="FFFFFF", line="D6DDE8", radius=True, shadow=True)
    add_text(slide, "输入", 1.1, 2.15, 0.8, 0.3, size=15, bold=True, color="2563EB")
    add_text(slide, "搜索词\n或 10 个竞品链接\n或 Amazon 前台页面", 1.1, 2.82, 2.15, 1.08, size=18, color="334155")
    add_chip(slide, "不需要导表", 1.1, 4.65, 1.25, fill="E7F7F4", color="0F766E")
    add_line(slide, 4.08, 3.8, 4.78, 3.8, color="F97316", width=2.2)
    add_rect(slide, 4.98, 1.75, 3.5, 4.25, fill="FFFFFF", line="D6DDE8", radius=True, shadow=True)
    add_text(slide, "Codex 先整理", 5.28, 2.03, 2.3, 0.3, size=15, bold=True, color="111827")
    items = ["价格带", "评论门槛", "主图打法", "卖点结构", "我们差距"]
    for i, item in enumerate(items):
        add_rect(slide, 5.32, 2.62 + i * 0.55, 2.55, 0.35, fill=["EAF1FF", "E7F7F4", "FFF3E8", "F3EAFE", "EAFBF4"][i], line=None, radius=True)
        add_text(slide, item, 5.48, 2.68 + i * 0.55, 2.1, 0.18, size=11.5, color="334155", bold=True, align="center")
    add_line(slide, 8.72, 3.8, 9.42, 3.8, color="F97316", width=2.2)
    add_rect(slide, 9.65, 1.9, 2.95, 3.95, fill="111827", line=None, radius=True, shadow=True)
    add_text(slide, "输出", 9.95, 2.18, 0.8, 0.3, size=15, bold=True, color="FFFFFF")
    add_text(slide, "这个市场靠什么卖？\n哪些点值得学？\n我们能不能打？\n下一步看什么证据？", 9.95, 2.82, 2.15, 1.5, size=15, color="E2E8F0")
    add_text(slide, "示例提示词：帮我看这个搜索词下的 10 个竞品，整理价格、评论、主图和卖点差异。", 1.0, 6.27, 11.4, 0.35, size=12.5, color="64748B", align="center")
    add_footer(slide, 6)
    return slide


def slide_conversion(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    add_background(slide)
    add_title(slide, "场景 2：转化分析，让它先找证据和反证", "这个比“帮我写复盘”更有吸引力：销售提出一个判断，Codex 去找支持它和推翻它的证据。", "场景")
    add_rect(slide, 0.85, 1.85, 3.15, 4.35, fill="FFFFFF", line="D6DDE8", radius=True, shadow=True)
    add_text(slide, "人的判断", 1.15, 2.14, 1.5, 0.3, size=15, bold=True, color="111827")
    add_text(slide, "我怀疑：\n这个 SKU 不是没流量，\n而是点击后不转化。", 1.15, 2.82, 2.35, 1.2, size=18, bold=True, color="2563EB")
    add_text(slide, "人负责提出假设，不负责把所有证据从头扒一遍。", 1.15, 5.1, 2.35, 0.48, size=11.5, color="64748B")
    evidence = [
        ("广告数据", "点击、CTR、花费、订单"),
        ("Listing 状态", "价格、评价、图片、卖点"),
        ("竞品对照", "同价位能否成交"),
        ("市场信号", "关键词需求、季节、价格带"),
    ]
    for i, (t, b) in enumerate(evidence):
        x = 4.55 + (i % 2) * 3.95
        y = 2.02 + (i // 2) * 1.58
        add_rect(slide, x, y, 3.3, 1.18, fill="FFFFFF", line="D6DDE8", radius=True, shadow=True)
        add_text(slide, t, x + 0.22, y + 0.2, 1.4, 0.25, size=14, bold=True, color=["2563EB", "0F766E", "F97316", "7C3AED"][i])
        add_text(slide, b, x + 0.22, y + 0.62, 2.75, 0.25, size=11.5, color="64748B")
    add_rect(slide, 4.55, 5.35, 7.25, 0.72, fill="111827", line=None, radius=True)
    add_text(slide, "输出不是一句“建议优化”，而是一张：支持证据 / 反证 / 还缺什么 / 下一步动作", 4.82, 5.55, 6.75, 0.22, size=12.5, color="FFFFFF", bold=True, align="center")
    add_footer(slide, 7)
    return slide


def slide_action_last_mile(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    add_background(slide)
    add_title(slide, "场景 3：最后一公里可以做，但要有闸门", "如果永远只分析，不动作，效率差距不够大。正确讲法是：先只读，成熟后按授权清单执行并回读。", "场景")
    steps = [
        ("1", "先生成清单", "按你的规则列出建议动作，不执行"),
        ("2", "人确认", "确认范围、强度、例外项和审批边界"),
        ("3", "授权执行", "只执行清单内动作，避免自由发挥"),
        ("4", "回读结果", "重新读取页面，确认动作是否真的落地"),
    ]
    x = 0.85
    for i, (num, title, body) in enumerate(steps):
        add_rect(slide, x + i * 3.05, 2.05, 2.45, 3.35, fill="FFFFFF", line="D6DDE8", radius=True, shadow=True)
        add_rect(slide, x + i * 3.05 + 0.28, 2.35, 0.64, 0.64, fill=["EAF1FF", "FFF3E8", "F3EAFE", "EAFBF4"][i], line=None, radius=True)
        add_text(slide, num, x + i * 3.05 + 0.28, 2.48, 0.64, 0.22, size=16, color=["2563EB", "F97316", "7C3AED", "0F766E"][i], bold=True, align="center", valign="middle")
        add_text(slide, title, x + i * 3.05 + 0.28, 3.28, 1.9, 0.34, size=15, bold=True, color="111827")
        add_text(slide, body, x + i * 3.05 + 0.28, 3.82, 1.88, 0.82, size=12, color="64748B")
        if i < 3:
            add_line(slide, x + i * 3.05 + 2.52, 3.72, x + i * 3.05 + 2.9, 3.72, color="CBD5E1", width=1.5)
    add_rect(slide, 1.62, 5.98, 10.1, 0.58, fill="111827", line=None, radius=True)
    add_text(slide, "会议上的一句话：可以做，但必须有授权、有清单、有回读。", 1.86, 6.12, 9.62, 0.22, size=13.5, color="FFFFFF", bold=True, align="center")
    add_footer(slide, 8)
    return slide


def slide_prompt_templates(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    add_background(slide)
    add_title(slide, "同事不用学复杂提示词，先从四句话开始", "这页的目的不是教 Prompt 工程，而是让大家有一个马上能复制的入口。", "低门槛")
    prompts = [
        ("只读查看", "只读帮我看这个 SKU 的广告页面，先不要做任何修改。"),
        ("竞品分析", "帮我看这个搜索词下的 10 个竞品，整理价格、评论、主图和卖点差异。"),
        ("找证据", "我怀疑这个 SKU 点击后不转化，帮我找支持和反证。"),
        ("动作清单", "按我的规则生成待确认动作清单，不要直接执行。"),
    ]
    for i, (title, prompt) in enumerate(prompts):
        x = 0.88 + (i % 2) * 6.15
        y = 1.9 + (i // 2) * 1.78
        add_rect(slide, x, y, 5.48, 1.38, fill="FFFFFF", line="D6DDE8", radius=True, shadow=True)
        add_chip(slide, title, x + 0.28, y + 0.25, 1.08, fill=["EAF1FF", "E7F7F4", "FFF3E8", "F3EAFE"][i], color=["2563EB", "0F766E", "F97316", "7C3AED"][i])
        add_text(slide, prompt, x + 0.3, y + 0.72, 4.85, 0.42, size=13.5, color="334155")
    add_rect(slide, 2.0, 5.96, 9.3, 0.48, fill="FFFFFF", line="D6DDE8", radius=True)
    add_text(slide, "关键加一句：先不要修改。这样新人也敢试，组长也能放心。", 2.2, 6.08, 8.9, 0.2, size=12.5, color="64748B", bold=True, align="center")
    add_footer(slide, 9)
    return slide


def slide_demo_plan(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    add_background(slide)
    add_title(slide, "现场演示不要大而全，只演示一条可见路径", "演示的目的不是证明它无所不能，而是证明：它真的能读你已登录 Chrome 里的业务页面。", "演示")
    add_rect(slide, 0.86, 1.92, 11.65, 3.0, fill="FFFFFF", line="D6DDE8", radius=True, shadow=True)
    flow = [
        ("一句话", "只读看这个 SKU"),
        ("打开页面", "Chrome 已登录"),
        ("读取字段", "广告 / 库存 / 状态"),
        ("生成初稿", "证据和待确认点"),
        ("人来判断", "改规则或确认"),
    ]
    for i, (t, b) in enumerate(flow):
        x = 1.25 + i * 2.2
        add_rect(slide, x, 2.48, 1.48, 0.88, fill=["EAF1FF", "E7F7F4", "FFF3E8", "F3EAFE", "EAFBF4"][i], line=None, radius=True)
        add_text(slide, t, x + 0.12, 2.62, 1.22, 0.22, size=12.5, bold=True, color=["2563EB", "0F766E", "F97316", "7C3AED", "0F766E"][i], align="center")
        add_text(slide, b, x + 0.12, 2.96, 1.22, 0.18, size=9.5, color="475569", align="center")
        if i < 4:
            add_line(slide, x + 1.52, 2.92, x + 2.0, 2.92, color="CBD5E1", width=1.4)
    add_text(slide, "可选演示素材：广告系统里一个 SKU，或 Amazon 前台一个搜索词。HUA0165 只作为技术已跑通证明，不把分享局限成单一案例。", 1.16, 4.08, 10.95, 0.42, size=12, color="64748B", align="center")
    add_rect(slide, 1.05, 5.42, 3.3, 0.64, fill="FEECEC", line=None, radius=True)
    add_text(slide, "不要承诺所有任务 3 分钟", 1.22, 5.58, 2.95, 0.18, size=12, color="B91C1C", bold=True, align="center")
    add_rect(slide, 5.02, 5.42, 3.3, 0.64, fill="EAFBF4", line=None, radius=True)
    add_text(slide, "承诺：先只读、可查验", 5.18, 5.58, 2.95, 0.18, size=12, color="047857", bold=True, align="center")
    add_rect(slide, 8.98, 5.42, 3.3, 0.64, fill="EAF1FF", line=None, radius=True)
    add_text(slide, "目标：让大家愿意试一次", 9.12, 5.58, 2.98, 0.18, size=12, color="1D4ED8", bold=True, align="center")
    add_footer(slide, 10)
    return slide


def slide_team_adoption(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    add_background(slide)
    add_title(slide, "组内推进：不强推工具，先沉淀共用方法", "组长在场时，这页很关键。重点不是“大家都马上用起来”，而是把团队重复经验变成可复用模板。", "落地")
    add_rect(slide, 0.92, 1.9, 3.4, 3.65, fill="FFFFFF", line="D6DDE8", radius=True, shadow=True)
    add_text(slide, "第 1 周", 1.2, 2.18, 1.0, 0.28, size=14.5, bold=True, color="2563EB")
    add_text(slide, "每个人只选一个\n最近重复做的小动作", 1.2, 2.95, 2.55, 0.8, size=18, color="334155")
    add_text(slide, "例如：竞品整理、转化证据、只读广告诊断。", 1.2, 4.52, 2.58, 0.34, size=11.5, color="64748B")
    add_rect(slide, 4.95, 1.9, 3.4, 3.65, fill="FFFFFF", line="D6DDE8", radius=True, shadow=True)
    add_text(slide, "第 2 步", 5.23, 2.18, 1.0, 0.28, size=14.5, bold=True, color="F97316")
    add_text(slide, "把好用的说法\n沉淀成 2-3 条模板", 5.23, 2.95, 2.55, 0.8, size=18, color="334155")
    add_text(slide, "不是每个人各玩各的，而是共用一套入口。", 5.23, 4.52, 2.58, 0.34, size=11.5, color="64748B")
    add_rect(slide, 8.98, 1.9, 3.4, 3.65, fill="111827", line=None, radius=True, shadow=True)
    add_text(slide, "成熟后", 9.26, 2.18, 1.0, 0.28, size=14.5, bold=True, color="FFFFFF")
    add_text(slide, "再讨论授权动作\n和固定流程", 9.26, 2.95, 2.55, 0.8, size=18, bold=True, color="FFFFFF")
    add_text(slide, "分析先跑通，修改后上车。", 9.26, 4.52, 2.58, 0.34, size=11.5, color="CBD5E1")
    add_text(slide, "这能避免一开始就把同事推到复杂自动化里。", 1.0, 6.1, 11.25, 0.32, size=12.5, color="64748B", align="center")
    add_footer(slide, 11)
    return slide


def slide_install(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    add_background(slide)
    add_title(slide, "安装不要一开始讲，放到“想试的人怎么开始”", "公司指南本身已经写得很全，现场只需要告诉大家：不用自己研究，卡住就找协助。", "会后")
    add_rect(slide, 0.85, 1.85, 4.05, 4.25, fill="FFFFFF", line="D6DDE8", radius=True, shadow=True)
    add_text(slide, "会后三步", 1.18, 2.15, 1.4, 0.32, size=16, bold=True, color="111827")
    install_steps = [
        "申请公司版 Codex",
        "按企业微信指南完成安装",
        "打开 Chrome，登录平时看的页面",
    ]
    for i, s in enumerate(install_steps):
        add_rect(slide, 1.2, 2.83 + i * 0.75, 0.42, 0.42, fill=["EAF1FF", "FFF3E8", "E7F7F4"][i], line=None, radius=True)
        add_text(slide, str(i + 1), 1.2, 2.91 + i * 0.75, 0.42, 0.16, size=10.5, color=["2563EB", "F97316", "0F766E"][i], bold=True, align="center")
        add_text(slide, s, 1.82, 2.84 + i * 0.75, 2.5, 0.24, size=13.5, color="334155")
    add_text(slide, "卡住不用硬扛：联系方燕 / 张文琦远程协助。", 1.18, 5.35, 3.15, 0.28, size=11.8, color="64748B")
    add_rect(slide, 5.48, 1.85, 3.05, 4.25, fill="FFFFFF", line="D6DDE8", radius=True, shadow=True)
    add_text(slide, "这里放二维码", 5.96, 3.02, 2.1, 0.3, size=15, color="94A3B8", bold=True, align="center")
    add_text(slide, "企业微信文档\n《Codex安装&使用教程》", 5.88, 3.5, 2.25, 0.55, size=12, color="64748B", align="center")
    add_rect(slide, 9.05, 1.85, 3.25, 4.25, fill="111827", line=None, radius=True, shadow=True)
    add_text(slide, "现场不要这么讲", 9.35, 2.18, 2.3, 0.32, size=15.5, bold=True, color="FFFFFF")
    add_text(slide, "“安装很简单，大家马上装。”", 9.35, 2.95, 2.35, 0.34, size=13, color="FCA5A5", bold=True)
    add_text(slide, "可以这么讲", 9.35, 3.72, 2.3, 0.32, size=15.5, bold=True, color="FFFFFF")
    add_text(slide, "“先看完效果。想试的人，会后按公司指南走；卡住有人远程帮。”", 9.35, 4.34, 2.38, 0.78, size=12.5, color="CBD5E1")
    add_footer(slide, 12)
    return slide


def slide_closing(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    add_background(slide)
    add_rect(slide, 0.75, 0.72, 11.83, 5.55, fill="111827", line=None, radius=True, shadow=True)
    add_text(slide, "今天先不用变成高手", 1.35, 1.35, 10.6, 0.62, size=31, bold=True, color="FFFFFF", align="center")
    add_text(slide, "找一个最近重复做的事，让 Codex 先帮你做一小步。", 1.7, 2.28, 9.95, 0.48, size=19, color="E2E8F0", align="center")
    add_rect(slide, 2.1, 3.32, 2.55, 1.24, fill="FFFFFF", line=None, radius=True)
    add_text(slide, "看页面", 2.35, 3.63, 2.0, 0.28, size=16, color="2563EB", bold=True, align="center")
    add_rect(slide, 5.38, 3.32, 2.55, 1.24, fill="FFFFFF", line=None, radius=True)
    add_text(slide, "找证据", 5.63, 3.63, 2.0, 0.28, size=16, color="F97316", bold=True, align="center")
    add_rect(slide, 8.66, 3.32, 2.55, 1.24, fill="FFFFFF", line=None, radius=True)
    add_text(slide, "列清单", 8.91, 3.63, 2.0, 0.28, size=16, color="10B981", bold=True, align="center")
    add_text(slide, "一句话收尾：重复动作越多，越应该让工具先跑第一遍。", 1.35, 5.3, 10.6, 0.32, size=13.5, color="CBD5E1", align="center")
    add_footer(slide, 13)
    return slide


SLIDE_NOTES = [
    ("封面", "先告诉大家：今天不是介绍一个新系统，而是让大家看到重复步骤可以被交出去一小步。"),
    ("开场效果", "开场不要讲概念，先放一个已登录页面只读分析的效果。注意不要承诺所有场景固定三分钟。"),
    ("不是聊天框", "把 ChatGPT 和 Codex 的差异讲清楚：一个主要靠人喂材料，一个可以跟随工作环境。"),
    ("Chrome Extension", "这页是王炸，但要稳：本人权限、可见页面、先只读、后确认。不要讲抓包或后台路径。"),
    ("重复事项", "给同事一个简单判断标准：只要每周做、步骤相似、有方法论、结果能查，就可以试。"),
    ("竞品市场", "强调不是看一个链接，而是批量看 10 个竞品或一个搜索词下的市场结构。"),
    ("转化证据", "把 AI 从文案工具变成证据工具：让它找支持证据和反证。"),
    ("最后一公里", "可以讲动作能力，但要加闸门：授权、清单、回读。"),
    ("四句入口", "让同事觉得很简单，不要讲 Prompt 工程。"),
    ("现场演示", "只演一条路径。HUA0165 只是证明能力，不把分享局限到一个 SKU 案例。"),
    ("组内推进", "组长在场时强调：不是强推大家自动化，是先沉淀可复制模板。"),
    ("安装会后讲", "安装放后面：先看效果，想试的人按公司指南，卡住找协助。"),
    ("收尾", "把行动要求降到最低：每个人找一个重复小动作试一次。"),
]


def write_notes():
    lines = ["# 销售 Codex 分享讨论稿 v1 - 讲稿要点", ""]
    for i, (title, note) in enumerate(SLIDE_NOTES, start=1):
        lines.append(f"## Slide {i:02d} - {title}")
        lines.append(note)
        lines.append("")
    NOTES_PATH.write_text("\n".join(lines), encoding="utf-8")


def build():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    pres = Presentation()
    pres.slide_width = Inches(W)
    pres.slide_height = Inches(H)
    pres.core_properties.title = "销售怎么开始用 Codex：从一个重复小动作开始"
    pres.core_properties.subject = "销售团队 Codex 分享培训讨论稿"
    pres.core_properties.author = "Codex"

    slide_cover(pres)
    slide_hook(pres)
    slide_not_chatgpt(pres)
    slide_chrome_extension(pres)
    slide_repetition_filter(pres)
    slide_competitor(pres)
    slide_conversion(pres)
    slide_action_last_mile(pres)
    slide_prompt_templates(pres)
    slide_demo_plan(pres)
    slide_team_adoption(pres)
    slide_install(pres)
    slide_closing(pres)

    pres.save(PPTX_PATH)
    write_notes()
    print(PPTX_PATH)
    print(NOTES_PATH)


if __name__ == "__main__":
    build()
